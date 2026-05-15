"""
Generates 'Day1_Part1_Intro_to_AI.pptx'
Audience: non-technical faculty. Goal: build intuition, not jargon.
Run:  py generate_intro_ai.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY      = RGBColor(0x0B, 0x2A, 0x4A)
TEAL      = RGBColor(0x00, 0x9B, 0x9E)
ACCENT    = RGBColor(0xF4, 0xA2, 0x61)
LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x2D, 0x3D)
GREY_TEXT = RGBColor(0x55, 0x66, 0x77)
SOFT_BLUE = RGBColor(0xE8, 0xF1, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=LIGHT_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    return bg

def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    return s

def add_round(slide, x, y, w, h, fill=WHITE, line=NAVY, adj=0.12):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = adj
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(1.25)
    return s

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def add_bullets(slide, x, y, w, h, items, size=16, color=DARK_TEXT, bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run(); r.text = f"{bullet}  {item}"
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def add_arrow(slide, x, y, w, h, color=TEAL):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a

def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(1.0), fill=NAVY)
    add_rect(slide, 0, Inches(1.0), SW, Inches(0.08), fill=TEAL)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.62), Inches(12), Inches(0.4),
                 subtitle, size=14, color=RGBColor(0xCF, 0xE3, 0xF5))

def footer(slide, page):
    add_text(slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.3),
             "Day 1 • Session 1 — Intro to AI", size=10, color=GREY_TEXT)
    add_text(slide, Inches(12.3), Inches(7.15), Inches(1), Inches(0.3),
             str(page), size=10, color=GREY_TEXT, align=PP_ALIGN.RIGHT)


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_rect(s, 0, Inches(3.2), SW, Inches(1.2), fill=TEAL)
add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "DAY 1  •  SESSION 1  of  4", size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(1.4),
         "Intro to AI", size=64, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.7),
         "What it is. How it thinks. Why now.",
         size=24, color=RGBColor(0xCF, 0xE3, 0xF5))
add_text(s, Inches(0.6), Inches(3.45), Inches(12), Inches(0.7),
         "An hour with no jargon, lots of intuition.",
         size=20, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
         "Two-Day Workshop on Agentic AI",
         size=14, color=ACCENT, bold=True)


# ---------- Slide 2: Agenda ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "What This Hour Looks Like")
items = [
    ("1", "What is AI, really?",         "Strip away the buzzwords"),
    ("2", "AI vs ML vs DL vs GenAI",     "One picture, no confusion"),
    ("3", "How does an LLM 'think'?",    "The next-word game"),
    ("4", "Why now? What changed?",      "The 3 reasons AI exploded"),
    ("5", "What AI is great at",         "And what it's terrible at"),
    ("6", "How we talk to AI",           "Prompts, tokens, context"),
]
for i, (n, t, d) in enumerate(items):
    row = i // 2; col = i % 2
    x = Inches(0.6 + col * 6.3)
    y = Inches(1.5 + row * 1.7)
    add_round(s, x, y, Inches(6.0), Inches(1.5), fill=WHITE, line=TEAL)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.35),
                              Inches(0.85), Inches(0.85))
    circ.fill.solid(); circ.fill.fore_color.rgb = NAVY
    circ.line.fill.background()
    add_text(s, x + Inches(0.25), y + Inches(0.38), Inches(0.85), Inches(0.8),
             n, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.3), y + Inches(0.3), Inches(4.5), Inches(0.5),
             t, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(1.3), y + Inches(0.78), Inches(4.5), Inches(0.6),
             d, size=13, color=GREY_TEXT)
footer(s, 2)


# ---------- Slide 3: A simple definition ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "So… what is AI?", "Forget the movies for a moment")

add_round(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.0),
          fill=NAVY, line=NAVY)
add_text(s, Inches(1.7), Inches(2.3), Inches(9.9), Inches(0.7),
         "AI is software that learns patterns from examples",
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.7), Inches(3.05), Inches(9.9), Inches(0.7),
         "instead of following hand-written rules.",
         size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.5),
         "A familiar analogy",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(2.0),
         "We don't teach a child what a 'cat' is by listing rules\n"
         "(“4 legs, fur, whiskers, pointy ears…”).\n"
         "We show them many cats. They figure out the pattern.\n"
         "AI learns the same way — just with millions of examples.",
         size=16, color=DARK_TEXT, align=PP_ALIGN.CENTER)
footer(s, 3)


# ---------- Slide 4: AI vs ML vs DL vs GenAI ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "AI, ML, Deep Learning, GenAI — Untangled",
       "They're not competitors. They're nested boxes.")

# Concentric rectangles
add_round(s, Inches(1.0), Inches(1.5), Inches(7.0), Inches(5.5),
          fill=SOFT_BLUE, line=NAVY, adj=0.04)
add_text(s, Inches(1.3), Inches(1.65), Inches(6.5), Inches(0.4),
         "Artificial Intelligence  (AI)", size=16, bold=True, color=NAVY)

add_round(s, Inches(1.4), Inches(2.3), Inches(6.2), Inches(4.5),
          fill=WHITE, line=TEAL, adj=0.05)
add_text(s, Inches(1.7), Inches(2.45), Inches(5.8), Inches(0.4),
         "Machine Learning  (ML)", size=15, bold=True, color=TEAL)

add_round(s, Inches(1.8), Inches(3.05), Inches(5.4), Inches(3.5),
          fill=LIGHT_BG, line=NAVY, adj=0.06)
add_text(s, Inches(2.1), Inches(3.2), Inches(5.0), Inches(0.4),
         "Deep Learning  (DL)", size=14, bold=True, color=NAVY)

add_round(s, Inches(2.2), Inches(3.8), Inches(4.6), Inches(2.6),
          fill=ACCENT, line=ACCENT, adj=0.08)
add_text(s, Inches(2.4), Inches(4.7), Inches(4.2), Inches(0.6),
         "Generative AI", size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(2.4), Inches(5.3), Inches(4.2), Inches(0.5),
         "(ChatGPT lives here)", size=13, color=WHITE,
         align=PP_ALIGN.CENTER)

# Right side explanations
defs = [
    ("AI",  "Any machine doing something 'smart'",   NAVY),
    ("ML",  "AI that learns from data",              TEAL),
    ("DL",  "ML using brain-inspired networks",      NAVY),
    ("GenAI","DL that creates new text/images/code", ACCENT),
]
for i, (term, desc, col) in enumerate(defs):
    y = Inches(1.7 + i * 1.3)
    add_round(s, Inches(8.6), y, Inches(4.2), Inches(1.1), fill=WHITE, line=col)
    add_text(s, Inches(8.8), y + Inches(0.15), Inches(3.8), Inches(0.4),
             term, size=18, bold=True, color=col)
    add_text(s, Inches(8.8), y + Inches(0.58), Inches(3.8), Inches(0.5),
             desc, size=13, color=DARK_TEXT)
footer(s, 4)


# ---------- Slide 5: How does an LLM think? ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "How Does ChatGPT 'Think'?",
       "Spoiler: it doesn't. It plays the world's best guessing game.")

add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "Everything an LLM does boils down to ONE skill:",
         size=18, color=DARK_TEXT)
add_text(s, Inches(0.7), Inches(1.95), Inches(12), Inches(0.7),
         "Predicting the next word.",
         size=32, bold=True, color=ACCENT)

# Example
add_text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(0.5),
         "You type:", size=14, color=GREY_TEXT)
add_round(s, Inches(0.7), Inches(3.4), Inches(11.9), Inches(0.7),
          fill=WHITE, line=NAVY)
add_text(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.5),
         "“The capital of France is ___”", size=18, bold=True, color=NAVY)

add_text(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.5),
         "The AI considers every possible next word and picks the most likely one:",
         size=13, color=GREY_TEXT)

# Probability chips
probs = [("Paris", "92%", ACCENT), ("Lyon", "3%", TEAL),
         ("a", "2%", GREY_TEXT), ("beautiful", "1%", GREY_TEXT)]
cx = Inches(0.7); cy = Inches(4.85)
for w, p, col in probs:
    add_round(s, cx, cy, Inches(2.5), Inches(0.9), fill=col, line=col)
    add_text(s, cx, cy + Inches(0.1), Inches(2.5), Inches(0.4),
             w, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, cx, cy + Inches(0.45), Inches(2.5), Inches(0.4),
             p, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    cx += Inches(2.7)

add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.6),
         "Then it adds that word, asks again, adds another, asks again…\n"
         "Billions of these tiny guesses → essays, code, lesson plans.",
         size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)
footer(s, 5)


# ---------- Slide 6: How did it get so smart? ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Then how is it so smart?",
       "Because it read… almost everything.")

stats = [
    ("📚", "Books", "Millions of them"),
    ("🌐", "Wikipedia", "Every article"),
    ("💻", "Code", "GitHub & more"),
    ("📰", "News & Web", "Years of articles"),
    ("💬", "Conversations", "Public forums, Q&A"),
]
for i, (ic, t, d) in enumerate(stats):
    x = Inches(0.5 + i * 2.55)
    y = Inches(1.5)
    add_round(s, x, y, Inches(2.4), Inches(2.4), fill=WHITE, line=TEAL)
    add_text(s, x, y + Inches(0.25), Inches(2.4), Inches(0.7),
             ic, size=36, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.15), Inches(2.4), Inches(0.5),
             t, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.65), Inches(2.4), Inches(0.5),
             d, size=12, color=GREY_TEXT, align=PP_ALIGN.CENTER)

add_round(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(2.4),
          fill=NAVY, line=NAVY)
add_text(s, Inches(1.0), Inches(4.55), Inches(11.5), Inches(0.6),
         "Imagine this:",
         size=18, bold=True, color=ACCENT)
add_text(s, Inches(1.0), Inches(5.1), Inches(11.5), Inches(1.5),
         "A student who has read every textbook, every research paper,\n"
         "every blog and every conversation on the internet —\n"
         "and remembers the patterns in all of it.\n"
         "That's roughly what a modern LLM is.",
         size=16, color=WHITE)
footer(s, 6)


# ---------- Slide 7: Why NOW? ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Why is AI Exploding NOW?",
       "The idea is old. Three things finally lined up.")

reasons = [
    ("📊", "Data",
     "The internet gave us oceans of text, images, and video.\n"
     "AI needs examples — and now we have billions."),
    ("⚡", "Compute",
     "GPUs designed for video games turned out to be perfect\n"
     "for training AI. Cloud made them rentable by the hour."),
    ("🧠", "Algorithms",
     "A 2017 breakthrough called the Transformer made it possible\n"
     "to train models on huge data efficiently."),
]
for i, (ic, t, d) in enumerate(reasons):
    y = Inches(1.5 + i * 1.85)
    add_round(s, Inches(0.6), y, Inches(12.1), Inches(1.65),
              fill=WHITE, line=TEAL)
    add_text(s, Inches(0.85), y + Inches(0.35), Inches(1.2), Inches(1.0),
             ic, size=44, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.2), y + Inches(0.2), Inches(10.0), Inches(0.5),
             t, size=22, bold=True, color=NAVY)
    add_text(s, Inches(2.2), y + Inches(0.75), Inches(10.0), Inches(0.9),
             d, size=13, color=DARK_TEXT)
footer(s, 7)


# ---------- Slide 8: What it's great at ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "What AI is Brilliant at  —  and Bad at",
       "Setting honest expectations")

# Good
add_round(s, Inches(0.5), Inches(1.4), Inches(6.1), Inches(5.4),
          fill=WHITE, line=TEAL)
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.1), Inches(0.7),
         fill=TEAL)
add_text(s, Inches(0.5), Inches(1.5), Inches(6.1), Inches(0.5),
         "✅  Great at", size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_bullets(s, Inches(0.8), Inches(2.3), Inches(5.6), Inches(4.5), [
    "Summarizing long documents",
    "Drafting emails, reports, lesson plans",
    "Explaining concepts in simple language",
    "Translating between languages",
    "Generating ideas & brainstorming",
    "Writing & explaining code",
], size=15)

# Bad
add_round(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.4),
          fill=WHITE, line=ACCENT)
add_rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.7),
         fill=ACCENT)
add_text(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.5),
         "⚠️  Watch out for", size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_bullets(s, Inches(7.1), Inches(2.3), Inches(5.6), Inches(4.5), [
    "Confidently making things up (hallucinations)",
    "Math & precise calculations",
    "Anything after its training cutoff date",
    "Your private/institutional data — it never saw it",
    "Remembering previous chats (no real memory)",
    "Being a source of truth without verification",
], size=15)
footer(s, 8)


# ---------- Slide 9: How we talk to AI - Prompts ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "How We Talk to AI  —  Prompts",
       "A prompt is just an instruction in plain English")

add_text(s, Inches(0.7), Inches(1.4), Inches(6.0), Inches(0.5),
         "Weak prompt", size=18, bold=True, color=GREY_TEXT)
add_round(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(1.2),
          fill=WHITE, line=GREY_TEXT)
add_text(s, Inches(0.9), Inches(2.15), Inches(5.6), Inches(0.8),
         "“Write something about AI.”",
         size=16, color=DARK_TEXT)

add_text(s, Inches(0.7), Inches(3.3), Inches(6.0), Inches(0.5),
         "Strong prompt", size=18, bold=True, color=TEAL)
add_round(s, Inches(0.7), Inches(3.8), Inches(6.0), Inches(2.8),
          fill=WHITE, line=TEAL)
add_text(s, Inches(0.9), Inches(3.95), Inches(5.6), Inches(2.6),
         "“You are a teacher explaining AI to first-year\n"
         "engineering students.\n"
         "Write a 200-word intro using a real-world\n"
         "analogy. Keep it friendly. Avoid jargon.”",
         size=14, color=DARK_TEXT)

# Right: the formula
add_round(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(5.4),
          fill=NAVY, line=NAVY)
add_text(s, Inches(7.4), Inches(1.6), Inches(5.2), Inches(0.6),
         "The 4-part recipe", size=22, bold=True, color=ACCENT)
recipe = [
    ("ROLE",    "Who should the AI act as?"),
    ("TASK",    "What exactly should it do?"),
    ("CONTEXT", "What background does it need?"),
    ("FORMAT",  "How should the answer look?"),
]
for i, (k, v) in enumerate(recipe):
    yy = Inches(2.5 + i * 1.0)
    add_text(s, Inches(7.4), yy, Inches(5.2), Inches(0.4),
             k, size=15, bold=True, color=ACCENT)
    add_text(s, Inches(7.4), yy + Inches(0.4), Inches(5.2), Inches(0.5),
             v, size=13, color=WHITE)
footer(s, 9)


# ---------- Slide 10: Tokens & Context ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Tokens, Context & Why They Matter",
       "The hidden 'units' the AI actually sees")

add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "AI doesn't read words. It reads ", size=16, color=DARK_TEXT)
add_text(s, Inches(3.45), Inches(1.4), Inches(8), Inches(0.5),
         "tokens.", size=16, bold=True, color=ACCENT)

# Token chips
add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.4),
         "“Faculty workshop on AI”  becomes:", size=14, color=GREY_TEXT)
chips = ["Faculty", "work", "shop", "on", "AI"]
cx = Inches(0.7); cy = Inches(2.5)
for w in chips:
    add_round(s, cx, cy, Inches(1.6), Inches(0.7), fill=TEAL, line=TEAL)
    add_text(s, cx, cy + Inches(0.15), Inches(1.6), Inches(0.4),
             w, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cx += Inches(1.8)

add_text(s, Inches(0.7), Inches(3.6), Inches(12), Inches(0.5),
         "Why faculty should care:",
         size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(4.1), Inches(12), Inches(2.5), [
    "More tokens in = more cost & slower response",
    "Every model has a 'context window' — a max number of tokens it can see at once",
    "If your document is too big, you must summarize or use RAG (post-lunch session!)",
    "This is the single biggest factor in real-world AI cost and design",
], size=15)
footer(s, 10)


# ---------- Slide 11: AI is already in your pocket ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "AI is Already in Your Pocket",
       "You and your students use it every single day — often without knowing")

cases = [
    ("🎬", "Netflix / YouTube",   "\u201cRecommended for you\u201d is AI"),
    ("🗺️", "Google Maps",         "Traffic + ETA prediction"),
    ("📷", "Phone Camera",        "Portrait mode, night mode, face unlock"),
    ("🛒", "Amazon / Flipkart",   "Product ranking & fraud checks"),
    ("💬", "WhatsApp / Gmail",    "Smart reply, spam filter, translate"),
    ("🎵", "Spotify",             "Discover Weekly is 100% AI"),
    ("🏦", "UPI / Banking",       "Real-time fraud detection"),
    ("🚗", "Uber / Ola",          "Surge pricing, driver matching"),
]
for i, (ic, t, d) in enumerate(cases):
    row = i // 4; col = i % 4
    x = Inches(0.4 + col * 3.18)
    y = Inches(1.4 + row * 2.6)
    add_round(s, x, y, Inches(3.05), Inches(2.4), fill=WHITE, line=TEAL)
    add_text(s, x, y + Inches(0.2), Inches(3.05), Inches(0.7),
             ic, size=32, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.0), Inches(3.05), Inches(0.5),
             t, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.5), Inches(3.05), Inches(0.8),
             d, size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
         "Your students are AI users already. Our job: turn them into AI builders.",
         size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
footer(s, 11)


# ---------- Slide 12: AI in the IT industry today ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "What the IT Industry is Actually Shipping (2024–26)",
       "The tools your students will use — or build — in their first job")

rows = [
    ("💻", "Coding",       "GitHub Copilot, Cursor, Claude Code, Devin",
                            "55% faster code completion (GitHub study)"),
    ("🎨", "Design & Media","Midjourney, DALL·E, Runway, ElevenLabs, Sora",
                            "Ads, posters, voice-overs, short films"),
    ("🔎", "Search & Q&A", "Perplexity, ChatGPT Search, Google AI Overviews",
                            "Search is being rewritten in real time"),
    ("🤖", "Agents",       "AutoGPT, LangChain, CrewAI, OpenAI Operator",
                            "AI that books, browses, and acts for you"),
    ("📞", "Customer Ops", "Intercom Fin, Zendesk AI, Salesforce Einstein",
                            "60–80% of L1 support handled by AI"),
    ("🧪", "Enterprise RAG","Glean, Notion AI, Microsoft Copilot for M365",
                            "Chat with your company's documents"),
]
for i, (ic, cat, tools, impact) in enumerate(rows):
    y = Inches(1.4 + i * 0.92)
    add_round(s, Inches(0.5), y, Inches(12.3), Inches(0.82),
              fill=WHITE, line=TEAL, adj=0.25)
    add_text(s, Inches(0.7), y + Inches(0.18), Inches(0.7), Inches(0.5),
             ic, size=24)
    add_text(s, Inches(1.5), y + Inches(0.12), Inches(2.3), Inches(0.35),
             cat, size=15, bold=True, color=NAVY)
    add_text(s, Inches(1.5), y + Inches(0.45), Inches(2.3), Inches(0.35),
             impact, size=10, color=ACCENT, bold=True)
    add_text(s, Inches(4.0), y + Inches(0.22), Inches(8.6), Inches(0.5),
             tools, size=13, color=DARK_TEXT)
footer(s, 12)


# ---------- Slide 13: Industries being transformed ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Industries Already Being Rewired by AI",
       "It's not just tech companies anymore")

industries = [
    ("🏥", "Healthcare",
     "AI reads X-rays & MRIs at radiologist-level.\n"
     "Drug discovery cut from years to months.\n"
     "Example: AlphaFold mapped 200M+ proteins."),
    ("💰", "Finance & Banking",
     "Fraud detection on every UPI transaction.\n"
     "AI traders, AI loan underwriting, AI KYC.\n"
     "JPMorgan saves 360k lawyer-hours/year."),
    ("🚗", "Automotive",
     "Tesla Autopilot, Waymo robotaxis.\n"
     "AI designs car parts lighter & stronger.\n"
     "Predictive maintenance for fleets."),
    ("🛍️", "Retail & E-commerce",
     "Personalised pricing & recommendations.\n"
     "AI demand forecasting, smart warehouses.\n"
     "Virtual try-on for clothes & makeup."),
    ("🏭", "Manufacturing",
     "Computer-vision defect detection.\n"
     "Digital twins for entire factories.\n"
     "Predictive maintenance saves crores."),
    ("📱", "Media & Marketing",
     "AI-generated ads, scripts, thumbnails.\n"
     "Hyper-personalised content at scale.\n"
     "Real-time dubbing into 30+ languages."),
]
for i, (ic, t, d) in enumerate(industries):
    row = i // 3; col = i % 3
    x = Inches(0.4 + col * 4.32)
    y = Inches(1.4 + row * 2.75)
    add_round(s, x, y, Inches(4.15), Inches(2.6), fill=WHITE, line=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(0.8), Inches(0.7),
             ic, size=30)
    add_text(s, x + Inches(1.0), y + Inches(0.2), Inches(3.0), Inches(0.5),
             t, size=16, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.95), Inches(3.7), Inches(1.6),
             d, size=11, color=DARK_TEXT)
footer(s, 13)


# ---------- Slide 14: Jobs your students should aim for ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "The Jobs Your Students Should Aim For",
       "What recruiters in 2026 are actually hiring — and what to teach")

jobs = [
    ("AI / ML Engineer",        "₹12–40 LPA", "Python, PyTorch, MLOps, system design"),
    ("GenAI / LLM Engineer",    "₹15–60 LPA", "Prompting, RAG, LangChain, vector DBs"),
    ("Data Engineer",           "₹10–30 LPA", "SQL, Spark, Airflow, cloud data lakes"),
    ("AI Product Manager",      "₹20–50 LPA", "Use-case scoping, evals, UX of AI apps"),
    ("MLOps / Platform",        "₹14–45 LPA", "Docker, K8s, model serving, monitoring"),
    ("AI Safety / Red Team",    "₹15–40 LPA", "Eval frameworks, jailbreak testing"),
]
add_text(s, Inches(0.5), Inches(1.3), Inches(7.5), Inches(0.4),
         "Hot roles  ·  India market", size=14, bold=True, color=NAVY)
for i, (role, pay, skills) in enumerate(jobs):
    y = Inches(1.75 + i * 0.78)
    add_round(s, Inches(0.5), y, Inches(7.5), Inches(0.7),
              fill=WHITE, line=TEAL, adj=0.3)
    add_text(s, Inches(0.7), y + Inches(0.08), Inches(4.0), Inches(0.35),
             role, size=13, bold=True, color=NAVY)
    add_text(s, Inches(0.7), y + Inches(0.38), Inches(4.0), Inches(0.3),
             skills, size=10, color=GREY_TEXT)
    add_text(s, Inches(4.8), y + Inches(0.2), Inches(2.6), Inches(0.4),
             pay, size=13, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

# Right panel: "What to add to your syllabus"
add_round(s, Inches(8.3), Inches(1.3), Inches(4.5), Inches(5.5),
          fill=NAVY, line=NAVY)
add_text(s, Inches(8.5), Inches(1.5), Inches(4.1), Inches(0.5),
         "What to teach NOW", size=18, bold=True, color=ACCENT)
teach = [
    "Python + Git fundamentals",
    "Calling LLM APIs (OpenAI / Gemini)",
    "Prompt engineering basics",
    "RAG: chat with your own data",
    "Vector databases (Chroma, FAISS)",
    "Agent frameworks (LangChain)",
    "Streamlit for quick demos",
    "Evaluating AI outputs critically",
]
for i, item in enumerate(teach):
    add_text(s, Inches(8.6), Inches(2.05 + i * 0.45), Inches(4.0), Inches(0.4),
             f"✓  {item}", size=12, color=WHITE)
footer(s, 14)


# ---------- Slide 15: Mindset shift ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "The Mindset Shift",
       "How to think about AI — not as magic, not as a threat")

add_round(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.3),
          fill=WHITE, line=GREY_TEXT)
add_rect(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.7), fill=GREY_TEXT)
add_text(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(0.5),
         "OLD WAY of thinking", size=18, bold=True, color=WHITE,  # (slide 15)
         align=PP_ALIGN.CENTER)
add_bullets(s, Inches(0.9), Inches(2.4), Inches(5.6), Inches(4.0), [
    "“AI will replace me.”",
    "“It's only for tech people.”",
    "“If it's wrong once, it's useless.”",
    "“I must learn coding first.”",
], size=15)

add_round(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3),
          fill=WHITE, line=ACCENT)
add_rect(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.7), fill=ACCENT)
add_text(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5),
         "NEW WAY of thinking", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_bullets(s, Inches(7.1), Inches(2.4), Inches(5.6), Inches(4.0), [
    "“AI is a faster intern I can supervise.”",
    "“Anyone who can write clearly can use it.”",
    "“I verify — but it saves me 70% of the draft time.”",
    "“I'll learn by using it on my real work.”",
], size=15)
footer(s, 15)


# ---------- Slide 16: Bridge to demo ----------
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_rect(s, 0, Inches(3.2), SW, Inches(0.08), fill=ACCENT)
add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
         "UP NEXT  —  Session 2", size=16, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(1.2),
         "Enough Theory.\nLet's See It Run.",
         size=52, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.6),
         "We'll take everything you just learned —",
         size=20, color=WHITE)
add_text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.6),
         "prompts, tokens, the next-word game —",
         size=20, color=WHITE)
add_text(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.6),
         "and watch a real AI Faculty Assistant work on student data.",
         size=20, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.5),
         "🎓  Live Demo — AI Faculty Assistant",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)


# ---------- Slide 17: Quick Q&A buffer ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Quick Questions Before We Demo?")
add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.7),
         "Anything unclear about:",
         size=20, bold=True, color=NAVY)
items = [
    "What AI actually is?",
    "How it 'thinks' (the next-word game)?",
    "Why it's brilliant at some things and bad at others?",
    "What prompts and tokens mean?",
]
add_bullets(s, Inches(1.2), Inches(3.0), Inches(11), Inches(4), items,
            size=20, bullet="›")
add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
         "Ask anything — no question is too basic.",
         size=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
footer(s, 17)


out = "Day1_Part1_Intro_to_AI.pptx"
prs.save(out)
print(f"Saved: {out}")
