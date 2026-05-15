"""
Generates 'AI_Faculty_Assistant_Presentation.pptx' — a faculty-friendly
walkthrough of the AI Faculty Assistant project.
Run:  py generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------- Theme ----------
NAVY       = RGBColor(0x0B, 0x2A, 0x4A)
TEAL       = RGBColor(0x00, 0x9B, 0x9E)
ACCENT     = RGBColor(0xF4, 0xA2, 0x61)
LIGHT_BG   = RGBColor(0xF5, 0xF7, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1F, 0x2D, 0x3D)
GREY_TEXT  = RGBColor(0x55, 0x66, 0x77)
GREEN_OK   = RGBColor(0x2E, 0xA0, 0x43)
RED_NO     = RGBColor(0xC2, 0x3B, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_bg(slide, color=LIGHT_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    return bg

def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s

def add_round(slide, x, y, w, h, fill=WHITE, line=NAVY):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.12
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(1.25)
    return s

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top  = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb

def add_bullets(slide, x, y, w, h, items, size=16, color=DARK_TEXT, bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run(); r.text = f"{bullet}  {item}"
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def add_arrow(slide, x, y, w, h, color=TEAL):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a

def add_down_arrow(slide, x, y, w, h, color=TEAL):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a

def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(1.0), fill=NAVY)
    add_rect(slide, 0, Inches(1.0), SW, Inches(0.08), fill=TEAL)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.62), Inches(12), Inches(0.4),
                 subtitle, size=14, color=RGBColor(0xCF, 0xE3, 0xF5))

def footer(slide, page):
    add_text(slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.3),
             "AI Faculty Assistant  •  Faculty Walkthrough",
             size=10, color=GREY_TEXT)
    add_text(slide, Inches(12.3), Inches(7.15), Inches(1), Inches(0.3),
             str(page), size=10, color=GREY_TEXT, align=PP_ALIGN.RIGHT)


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, Inches(3.2), SW, Inches(1.2), fill=TEAL)
add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "DAY 1  •  SESSION 2  of  4",
         size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.2),
         "🎓  AI Faculty Assistant",
         size=54, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.7),
         "Turning Student Data into Instant, Intelligent Insights",
         size=22, color=RGBColor(0xCF, 0xE3, 0xF5))
add_text(s, Inches(0.6), Inches(3.45), Inches(12), Inches(0.7),
         "From Theory  →  to a Working AI App",
         size=20, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.4),
         "Streamlit  •  Pandas  •  OpenAI GPT-4o-mini",
         size=16, color=ACCENT, bold=True)
add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
         "Two-Day Workshop on Agentic AI",
         size=14, color=WHITE)

# ---------- Slide 1.5: Workshop Context ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Where We Are in the Workshop",
       "We just covered the theory — now let's see it run")

sessions = [
    ("Session 1", "Intro to AI",           "What it is, how it thinks, why now", False),
    ("Session 2", "Live Demo — This App",  "Theory → working AI app",            True),
    ("Lunch",     "Break",                  "",                                   False),
    ("Session 3", "RAG Architecture",      "Scaling AI to real data",            False),
    ("Day 2",     "Agentic AI",            "When AI starts to act",              False),
]
y = Inches(1.5)
for i, (tag, name, desc, active) in enumerate(sessions):
    yy = y + Inches(i * 1.05)
    fill = ACCENT if active else WHITE
    line = ACCENT if active else TEAL
    add_round(s, Inches(0.6), yy, Inches(12.1), Inches(0.9), fill=fill, line=line)
    tag_col = WHITE if active else NAVY
    name_col = WHITE if active else NAVY
    desc_col = WHITE if active else GREY_TEXT
    add_text(s, Inches(0.85), yy + Inches(0.25), Inches(1.8), Inches(0.5),
             tag, size=14, bold=True, color=tag_col)
    add_text(s, Inches(2.8), yy + Inches(0.13), Inches(4.5), Inches(0.5),
             name, size=18, bold=True, color=name_col)
    add_text(s, Inches(2.8), yy + Inches(0.5), Inches(9.5), Inches(0.4),
             desc, size=12, color=desc_col)
    if active:
        add_text(s, Inches(11.3), yy + Inches(0.25), Inches(1.3), Inches(0.5),
                 "◀  YOU ARE HERE", size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT)
footer(s, 2)


# ---------- Slide 2: Agenda ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "What We'll Cover Today")
items = [
    ("1", "The Problem",            "Why faculty need an AI assistant"),
    ("2", "What This App Does",     "A quick tour of the live demo"),
    ("3", "How It Works",           "End-to-end flow in plain English"),
    ("4", "The Tech Stack",         "Streamlit, Pandas, OpenAI"),
    ("5", "Tokens & Cost",          "Why model choice matters"),
    ("6", "Today vs. Production",   "From simple prompt to RAG"),
    ("7", "What's Next",            "From assistant to Agentic AI"),
]
for i, (n, t, d) in enumerate(items):
    row = i // 2; col = i % 2
    x = Inches(0.6 + col * 6.3)
    y = Inches(1.5 + row * 1.35)
    add_round(s, x, y, Inches(6.0), Inches(1.15), fill=WHITE, line=TEAL)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2),
                              Inches(0.75), Inches(0.75))
    circ.fill.solid(); circ.fill.fore_color.rgb = NAVY; circ.line.fill.background()
    add_text(s, x + Inches(0.2), y + Inches(0.22), Inches(0.75), Inches(0.7),
             n, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.1), y + Inches(0.18), Inches(4.8), Inches(0.5),
             t, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(1.1), y + Inches(0.6),  Inches(4.8), Inches(0.5),
             d, size=13, color=GREY_TEXT)
footer(s, 2)


# ---------- Slide 3: The Problem ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "The Problem", "Faculty spend hours on questions data could answer in seconds")
pains = [
    "Hundreds of student records in Excel sheets",
    "Repetitive questions: top performers? weak students? attendance dips?",
    "Manual filtering, sorting, and reporting eats teaching time",
    "Insights are reactive, not proactive",
]
add_text(s, Inches(0.7), Inches(1.4), Inches(6.0), Inches(0.5),
         "Today's reality", size=22, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(2.0), Inches(6.0), Inches(4.5), pains, size=17)

# Right side: "What if..." card
add_round(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(5.0), fill=NAVY, line=NAVY)
add_text(s, Inches(7.4), Inches(1.6), Inches(5.2), Inches(0.6),
         "What if you could just… ask?", size=22, bold=True, color=ACCENT)
asks = [
    "“Who are my top 5 students this term?”",
    "“List students with attendance below 75%.”",
    "“Suggest mentoring actions for weak students.”",
    "“Give me a summary report for the HOD.”",
]
add_bullets(s, Inches(7.4), Inches(2.5), Inches(5.2), Inches(3.5), asks,
            size=15, color=WHITE, bullet="›")
footer(s, 3)


# ---------- Slide 4: What the App Does ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "What This App Does", "A lightweight, conversational layer over your Excel data")
features = [
    ("📤", "Upload Excel",        "Faculty uploads any .xlsx student file"),
    ("👀", "Preview Data",        "Instant tabular view inside the browser"),
    ("💬", "Ask Anything",        "Free-form questions about the data"),
    ("🧠", "AI Reasoning",        "GPT-4o-mini interprets and answers"),
    ("📊", "One-Click Insights",  "Top, weak, attendance, recommendations"),
    ("⚡", "Zero Setup for User", "Just a web link — no Excel formulas"),
]
for i, (ic, t, d) in enumerate(features):
    row = i // 3; col = i % 3
    x = Inches(0.5 + col * 4.25)
    y = Inches(1.5 + row * 2.6)
    add_round(s, x, y, Inches(4.0), Inches(2.3), fill=WHITE, line=TEAL)
    add_text(s, x, y + Inches(0.2), Inches(4.0), Inches(0.8),
             ic, size=36, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.05), Inches(4.0), Inches(0.5),
             t, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.25), y + Inches(1.55), Inches(3.5), Inches(0.7),
             d, size=12, color=GREY_TEXT, align=PP_ALIGN.CENTER)
footer(s, 4)


# ---------- Slide 5: How It Works (Flow) ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "How It Works", "End-to-end flow, in five simple steps")

steps = [
    ("Faculty Uploads Excel",   "📤"),
    ("Python Reads the Data",   "🐍"),
    ("Data + Question → LLM",   "📡"),
    ("LLM Understands Context", "🧠"),
    ("AI Generates Response",   "💡"),
]
box_w, box_h = Inches(2.25), Inches(1.7)
gap = Inches(0.13)
total_w = box_w * 5 + gap * 4
start_x = (SW - total_w) / 2
y = Inches(2.6)
for i, (t, ic) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    add_round(s, x, y, box_w, box_h, fill=WHITE, line=NAVY)
    add_text(s, x, y + Inches(0.15), box_w, Inches(0.5),
             ic, size=28, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(0.75), box_w - Inches(0.2), Inches(0.9),
             t, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    if i < 4:
        ax = x + box_w + Inches(-0.05)
        add_arrow(s, ax, y + Inches(0.7), gap + Inches(0.1), Inches(0.3))

add_text(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.5),
         "Plain-English summary",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(5.3), Inches(10.3), Inches(1.5),
         "We hand the AI two things — your spreadsheet and your question — "
         "and it returns a written answer. No SQL, no formulas, no code.",
         size=15, color=DARK_TEXT, align=PP_ALIGN.CENTER)
footer(s, 5)


# ---------- Slide 6: Tech Stack ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "The Tech Stack", "Three small pieces, one powerful experience")
techs = [
    ("Streamlit", "Frontend",
     "Build AI web apps in pure Python — no HTML, CSS, or JS required.",
     "Why: minimal frontend work, fast to demo."),
    ("Pandas", "Data Layer",
     "Reads and processes structured Excel data into a DataFrame.",
     "Why: industry standard for tabular data in Python."),
    ("OpenAI API", "Reasoning Engine",
     "GPT-4o-mini interprets the question and the data, then responds.",
     "Why: smaller model = lower cost, faster, ideal for a demo."),
]
for i, (name, role, desc, why) in enumerate(techs):
    x = Inches(0.5 + i * 4.27)
    y = Inches(1.5)
    add_round(s, x, y, Inches(4.05), Inches(5.2), fill=WHITE, line=TEAL)
    add_rect(s, x, y, Inches(4.05), Inches(0.9), fill=NAVY)
    add_text(s, x, y + Inches(0.12), Inches(4.05), Inches(0.5),
             name, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.55), Inches(4.05), Inches(0.4),
             role, size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.25), y + Inches(1.2), Inches(3.6), Inches(2.0),
             desc, size=14, color=DARK_TEXT)
    add_text(s, x + Inches(0.25), y + Inches(3.6), Inches(3.6), Inches(1.5),
             why, size=13, color=TEAL, bold=True)
footer(s, 6)


# ---------- Slide 7: Tokens Explained ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "What Are Tokens?", "The currency of every AI conversation")

add_text(s, Inches(0.7), Inches(1.4), Inches(6.0), Inches(0.5),
         "Every sentence is broken into tokens", size=20, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(2.0), Inches(6.0), Inches(0.5),
         "Example:", size=16, bold=True, color=DARK_TEXT)
add_text(s, Inches(0.7), Inches(2.4), Inches(6.0), Inches(0.6),
         "“Explain AI”", size=22, bold=True, color=DARK_TEXT)
add_text(s, Inches(0.7), Inches(3.05), Inches(6.0), Inches(0.4),
         "is seen by the model as:", size=14, color=GREY_TEXT)

# token chips
chips = ["Explain", "AI"]
cx = Inches(0.7); cy = Inches(3.55)
for w in chips:
    chip = add_round(s, cx, cy, Inches(1.5), Inches(0.7), fill=ACCENT, line=ACCENT)
    add_text(s, cx, cy + Inches(0.15), Inches(1.5), Inches(0.5),
             w, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cx += Inches(1.7)

add_text(s, Inches(0.7), Inches(4.6), Inches(6.0), Inches(2.5),
         "Every request consumes tokens from:\n"
         "   •  the prompt you write\n"
         "   •  the Excel data attached\n"
         "   •  the AI's response",
         size=15, color=DARK_TEXT)

# Right card: cost reality
add_round(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(5.0), fill=NAVY, line=NAVY)
add_text(s, Inches(7.4), Inches(1.6), Inches(5.2), Inches(0.6),
         "More rows  →  More tokens  →  More cost",
         size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(7.4), Inches(2.5), Inches(5.2), Inches(3.5), [
    "A 1,000-row Excel = thousands of tokens, every single query",
    "We chose GPT-4o-mini → cheaper, faster, perfect for demo scale",
    "In real enterprises, token optimization = scalability + cost control",
    "This is exactly why RAG architectures exist (next slide)",
], size=14, color=WHITE, bullet="›")
footer(s, 7)


# ---------- Slide 8: Today vs. Production (RAG) ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Today vs. Production", "What we built — and what it grows into")

# Left: Today
add_round(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.4), fill=WHITE, line=TEAL)
add_text(s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.5),
         "✅  Today — Lightweight LLM Assistant",
         size=18, bold=True, color=GREEN_OK)
add_text(s, Inches(0.7), Inches(2.1), Inches(5.6), Inches(0.4),
         "Excel  →  Prompt  →  LLM",
         size=16, bold=True, color=NAVY)
flow1 = ["Excel", "Prompt", "LLM"]
fx = Inches(0.7); fy = Inches(2.7)
for i, w in enumerate(flow1):
    add_round(s, fx, fy, Inches(1.3), Inches(0.6), fill=LIGHT_BG, line=NAVY)
    add_text(s, fx, fy + Inches(0.13), Inches(1.3), Inches(0.4),
             w, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    if i < 2:
        add_arrow(s, fx + Inches(1.32), fy + Inches(0.18), Inches(0.35), Inches(0.25))
    fx += Inches(1.7)
add_bullets(s, Inches(0.7), Inches(3.7), Inches(5.6), Inches(2.8), [
    "Best for small datasets",
    "Quick to build and demo",
    "Entire data sent every time",
    "Good for POCs, classroom demos",
], size=14)

# Right: Production RAG
add_round(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.4), fill=WHITE, line=ACCENT)
add_text(s, Inches(7.0), Inches(1.55), Inches(5.6), Inches(0.5),
         "🚀  Production — RAG Architecture",
         size=18, bold=True, color=ACCENT)
add_text(s, Inches(7.0), Inches(2.1), Inches(5.6), Inches(0.4),
         "Docs → Embeddings → Vector DB → Retrieval → LLM",
         size=14, bold=True, color=NAVY)
flow2 = ["Docs", "Embed", "Vector DB", "Retrieve", "LLM"]
fx = Inches(7.0); fy = Inches(2.7)
for i, w in enumerate(flow2):
    add_round(s, fx, fy, Inches(1.0), Inches(0.6), fill=LIGHT_BG, line=ACCENT)
    add_text(s, fx, fy + Inches(0.13), Inches(1.0), Inches(0.4),
             w, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    if i < 4:
        add_arrow(s, fx + Inches(1.02), fy + Inches(0.18),
                  Inches(0.13), Inches(0.25), color=ACCENT)
    fx += Inches(1.15)
add_bullets(s, Inches(7.0), Inches(3.7), Inches(5.6), Inches(2.8), [
    "Handles thousands of PDFs, papers, policies",
    "Sends only the relevant snippets to the LLM",
    "Massive cost savings at scale",
    "Foundation for enterprise-grade AI",
], size=14)
footer(s, 8)


# ---------- Slide 9: From Assistant to Agentic AI ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "What's Next: Agentic AI",
       "Today this assistant answers questions. Tomorrow it acts on them.")

add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.5),
         "Today  —  Question & Answer",
         size=18, bold=True, color=NAVY)
add_round(s, Inches(0.7), Inches(2.05), Inches(12), Inches(0.7), fill=WHITE, line=TEAL)
add_text(s, Inches(0.9), Inches(2.18), Inches(11.6), Inches(0.5),
         "Faculty asks → AI replies. The human takes the next action.",
         size=14, color=DARK_TEXT)

add_down_arrow(s, Inches(6.4), Inches(2.85), Inches(0.5), Inches(0.5))

add_text(s, Inches(0.7), Inches(3.45), Inches(12), Inches(0.5),
         "Tomorrow  —  Agentic AI",
         size=18, bold=True, color=ACCENT)
agents = [
    ("🔍",  "Identify weak students automatically"),
    ("✉️",  "Draft & send mentoring emails to parents"),
    ("📅",  "Schedule mentoring sessions on the calendar"),
    ("📑",  "Generate weekly faculty & HOD reports"),
]
for i, (ic, t) in enumerate(agents):
    row = i // 2; col = i % 2
    x = Inches(0.7 + col * 6.1)
    y = Inches(4.1 + row * 1.2)
    add_round(s, x, y, Inches(5.9), Inches(1.0), fill=WHITE, line=ACCENT)
    add_text(s, x + Inches(0.15), y + Inches(0.2), Inches(0.7), Inches(0.6),
             ic, size=24, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.95), y + Inches(0.28), Inches(4.8), Inches(0.5),
             t, size=15, bold=True, color=NAVY)

add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
         "“That's where Agentic AI begins — AI that doesn't just think, it acts.”",
         size=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
footer(s, 9)


# ---------- Slide 10: The Big Question (transition) ----------
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_rect(s, 0, Inches(2.8), SW, Inches(0.08), fill=ACCENT)
add_text(s, Inches(0.6), Inches(0.8), Inches(12), Inches(1.0),
         "What we built works beautifully…",
         size=32, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1.0),
         "…for small datasets.",
         size=26, color=ACCENT, bold=True)
add_text(s, Inches(0.6), Inches(3.1), Inches(12), Inches(0.6),
         "But what happens when we have…",
         size=20, color=WHITE)
qs = [
    "thousands of PDFs?",
    "research papers?",
    "policies & regulations?",
    "years of student records?",
]
for i, q in enumerate(qs):
    x = Inches(0.6 + (i % 2) * 6.3)
    y = Inches(3.9 + (i // 2) * 1.0)
    add_round(s, x, y, Inches(6.0), Inches(0.8), fill=WHITE, line=ACCENT)
    add_text(s, x + Inches(0.25), y + Inches(0.2), Inches(5.7), Inches(0.5),
             "›  " + q, size=18, bold=True, color=NAVY)
add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.6),
         "Can we still send everything to GPT every single time?",
         size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ---------- Slide 11: Key Takeaways ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Key Takeaways")
points = [
    ("1", "AI is now a tool faculty can use directly — no coding required."),
    ("2", "Streamlit + Pandas + OpenAI is enough to build a real assistant."),
    ("3", "Tokens = cost. Choosing the right model matters."),
    ("4", "Today's prompt-based app evolves into RAG for enterprise scale."),
    ("5", "The next frontier is Agentic AI — assistants that take action."),
]
for i, (n, t) in enumerate(points):
    y = Inches(1.6 + i * 1.0)
    add_round(s, Inches(0.6), y, Inches(12.1), Inches(0.85), fill=WHITE, line=TEAL)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.13),
                              Inches(0.6), Inches(0.6))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT; circ.line.fill.background()
    add_text(s, Inches(0.8), y + Inches(0.18), Inches(0.6), Inches(0.5),
             n, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.6), y + Inches(0.22), Inches(11.0), Inches(0.5),
             t, size=16, color=DARK_TEXT)
footer(s, 11)


# ---------- Slide 12: Thank You ----------
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_rect(s, 0, Inches(3.3), SW, Inches(0.08), fill=ACCENT)
add_text(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.2),
         "Thank You",
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(3.6), Inches(12.1), Inches(0.6),
         "Questions, ideas, and use-cases welcome.",
         size=20, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.5),
         "🎓  AI Faculty Assistant   •   Live Demo Ready",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)


out = "AI_Faculty_Assistant_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
