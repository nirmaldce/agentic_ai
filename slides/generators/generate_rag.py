"""
Generates 'Day1_Part3_RAG.pptx'
Audience: non-technical faculty. Conceptual with a light technical walk.
Keeps it lively with the Library / Librarian / Open-Book Exam analogy.
Run:  py generate_rag.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY      = RGBColor(0x0B, 0x2A, 0x4A)
TEAL      = RGBColor(0x00, 0x9B, 0x9E)
ACCENT    = RGBColor(0xF4, 0xA2, 0x61)
PURPLE    = RGBColor(0x7A, 0x4E, 0xAB)
LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x2D, 0x3D)
GREY_TEXT = RGBColor(0x55, 0x66, 0x77)
SOFT_BLUE = RGBColor(0xE8, 0xF1, 0xF8)
GREEN_OK  = RGBColor(0x2E, 0xA0, 0x43)
RED_NO    = RGBColor(0xC2, 0x3B, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=LIGHT_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    return bg

def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    return s

def add_round(slide, x, y, w, h, fill=WHITE, line=NAVY, adj=0.12):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = adj
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(1.25)
    return s

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def add_bullets(slide, x, y, w, h, items, size=16, color=DARK_TEXT, bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run(); r.text = f"{bullet}  {item}"
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def add_arrow(slide, x, y, w, h, color=TEAL):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a

def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(1.0), fill=NAVY)
    add_rect(slide, 0, Inches(1.0), SW, Inches(0.08), fill=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.62), Inches(12), Inches(0.4),
                 subtitle, size=14, color=RGBColor(0xCF, 0xE3, 0xF5))

def footer(slide, page):
    add_text(slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.3),
             "Day 1 • Session 3 — RAG Architecture", size=10, color=GREY_TEXT)
    add_text(slide, Inches(12.3), Inches(7.15), Inches(1), Inches(0.3),
             str(page), size=10, color=GREY_TEXT, align=PP_ALIGN.RIGHT)


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_rect(s, 0, Inches(3.2), SW, Inches(1.2), fill=ACCENT)
add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "DAY 1  •  SESSION 3  of  4   —   POST LUNCH",
         size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(1.4),
         "RAG", size=120, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(2.7), Inches(12), Inches(0.7),
         "Retrieval-Augmented Generation",
         size=30, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.5), Inches(12), Inches(0.7),
         "Giving AI access to YOUR knowledge — safely, cheaply, accurately.",
         size=20, color=RGBColor(0xCF, 0xE3, 0xF5))
add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
         "Two-Day Workshop on Agentic AI",
         size=14, color=ACCENT, bold=True)


# ---------- Slide 2: Recap + transition ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Before Lunch, We Saw…",
       "Excel → Prompt → LLM. Beautiful. But limited.")

add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "The demo worked because the dataset was small.",
         size=18, color=DARK_TEXT)

add_text(s, Inches(0.7), Inches(2.1), Inches(12), Inches(0.5),
         "But what if your knowledge base is…",
         size=20, bold=True, color=NAVY)

boxes = [
    ("📚", "10,000 PDFs"),
    ("📜", "20 years of policies"),
    ("🎓", "Every research paper in your dept"),
    ("📋", "All NAAC / NBA documents"),
]
for i, (ic, t) in enumerate(boxes):
    x = Inches(0.6 + (i % 2) * 6.3)
    y = Inches(3.0 + (i // 2) * 1.6)
    add_round(s, x, y, Inches(6.0), Inches(1.4), fill=WHITE, line=TEAL)
    add_text(s, x + Inches(0.3), y + Inches(0.35), Inches(1.0), Inches(0.8),
             ic, size=36, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.4), y + Inches(0.45), Inches(4.4), Inches(0.6),
             t, size=20, bold=True, color=NAVY)

add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
         "Can we paste all of that into every prompt?  No. So what then?",
         size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
footer(s, 2)


# ---------- Slide 3: The problem with stuffing ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Why We Can't Just 'Send Everything'",
       "Three hard limits of the simple approach")

probs = [
    ("💰", "It's expensive",
     "Every token costs money. 10,000 PDFs = millions of tokens per query."),
    ("📏", "It doesn't fit",
     "Every LLM has a context window. Even huge ones can't hold a library."),
    ("🎯", "It's distracting",
     "More irrelevant text in → worse, less focused answers out."),
]
for i, (ic, t, d) in enumerate(probs):
    x = Inches(0.5 + i * 4.27)
    y = Inches(1.5)
    add_round(s, x, y, Inches(4.05), Inches(5.0), fill=WHITE, line=ACCENT)
    add_text(s, x, y + Inches(0.4), Inches(4.05), Inches(1.0),
             ic, size=56, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.6), Inches(4.05), Inches(0.5),
             t, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + Inches(2.3), Inches(3.4), Inches(2.5),
             d, size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
         "The fix?  Don't send everything. Send only what's relevant.",
         size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
footer(s, 3)


# ---------- Slide 4: THE Big Analogy ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "The Open-Book Exam Analogy",
       "This is what RAG really is.")

add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "Imagine two students taking the same exam:",
         size=18, bold=True, color=NAVY)

# Student A
add_round(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.7),
          fill=WHITE, line=GREY_TEXT)
add_rect(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(0.8), fill=GREY_TEXT)
add_text(s, Inches(0.6), Inches(2.25), Inches(6.0), Inches(0.5),
         "Student A — Closed Book", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.85), Inches(3.1), Inches(5.6), Inches(0.5),
         "Just an LLM",
         size=14, bold=True, color=GREY_TEXT)
add_bullets(s, Inches(0.85), Inches(3.6), Inches(5.6), Inches(3.0), [
    "Relies only on what it memorized",
    "Doesn't know YOUR data",
    "May confidently make things up",
    "Cutoff date = stale info",
], size=14)

# Student B
add_round(s, Inches(6.8), Inches(2.1), Inches(6.0), Inches(4.7),
          fill=WHITE, line=ACCENT)
add_rect(s, Inches(6.8), Inches(2.1), Inches(6.0), Inches(0.8), fill=ACCENT)
add_text(s, Inches(6.8), Inches(2.25), Inches(6.0), Inches(0.5),
         "Student B — Open Book (RAG)", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(7.05), Inches(3.1), Inches(5.6), Inches(0.5),
         "LLM + the right pages, fetched live",
         size=14, bold=True, color=ACCENT)
add_bullets(s, Inches(7.05), Inches(3.6), Inches(5.6), Inches(3.0), [
    "Looks up the exact relevant pages first",
    "Grounds every answer in your documents",
    "Cites sources you can verify",
    "Stays up-to-date as you add new docs",
], size=14)
footer(s, 4)


# ---------- Slide 5: RAG in one sentence ----------
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_rect(s, 0, Inches(3.5), SW, Inches(0.08), fill=ACCENT)
add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
         "RAG, in one sentence:", size=22, color=ACCENT, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.5),
         "Find the right snippets,",
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(1.5),
         "hand them to the LLM,",
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.3), Inches(12), Inches(1.5),
         "ask the question.",
         size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ---------- Slide 6: The 5-step pipeline ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "The RAG Pipeline  —  5 Steps",
       "Two phases: prepare once, then answer forever")

# Phase 1
add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.4),
         "PHASE 1  —  Prepare your knowledge (done once, or whenever docs change)",
         size=13, bold=True, color=TEAL)
steps1 = [("📄", "Ingest", "Load PDFs,\ndocs, web"),
          ("✂️", "Chunk", "Split into\nsmall pieces"),
          ("🔢", "Embed", "Turn each piece\ninto numbers"),
          ("🗄️", "Store", "Save in a\nVector DB")]
box_w = Inches(2.6); gap = Inches(0.25)
total = box_w * 4 + gap * 3
sx = (SW - total) / 2; y = Inches(1.75)
for i, (ic, t, d) in enumerate(steps1):
    x = sx + i * (box_w + gap)
    add_round(s, x, y, box_w, Inches(1.6), fill=WHITE, line=TEAL)
    add_text(s, x, y + Inches(0.1), box_w, Inches(0.5),
             ic, size=24, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.5), box_w, Inches(0.4),
             t, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.9), box_w, Inches(0.7),
             d, size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)
    if i < 3:
        add_arrow(s, x + box_w - Inches(0.05), y + Inches(0.6),
                  gap + Inches(0.1), Inches(0.25), color=TEAL)

# Phase 2
add_text(s, Inches(0.7), Inches(3.85), Inches(12), Inches(0.4),
         "PHASE 2  —  Answer the question (every time the user asks)",
         size=13, bold=True, color=ACCENT)
steps2 = [("❓", "User asks", "“What's the\nleave policy?”"),
          ("🔍", "Retrieve", "Find the most\nrelevant chunks"),
          ("📦", "Build prompt", "Question + chunks\n→ context"),
          ("🤖", "LLM answers", "Grounded in\nyour docs")]
y = Inches(4.3)
for i, (ic, t, d) in enumerate(steps2):
    x = sx + i * (box_w + gap)
    add_round(s, x, y, box_w, Inches(1.6), fill=WHITE, line=ACCENT)
    add_text(s, x, y + Inches(0.1), box_w, Inches(0.5),
             ic, size=24, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.5), box_w, Inches(0.4),
             t, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.9), box_w, Inches(0.7),
             d, size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)
    if i < 3:
        add_arrow(s, x + box_w - Inches(0.05), y + Inches(0.6),
                  gap + Inches(0.1), Inches(0.25), color=ACCENT)

add_text(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.4),
         "We'll walk through each piece next — gently.",
         size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
footer(s, 6)


# ---------- Slide 7: Chunking ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Step 1  —  Chunking",
       "Why we slice documents into bite-sized pieces")

add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "A 200-page PDF is too big to retrieve as one unit.",
         size=18, color=DARK_TEXT)

# Big doc -> chunks visual
add_round(s, Inches(0.7), Inches(2.2), Inches(2.5), Inches(3.2),
          fill=WHITE, line=NAVY)
add_text(s, Inches(0.7), Inches(2.4), Inches(2.5), Inches(0.4),
         "200-page PDF", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
for i in range(6):
    add_rect(s, Inches(1.0), Inches(2.9 + i * 0.35), Inches(1.9),
             Inches(0.25), fill=SOFT_BLUE)

add_arrow(s, Inches(3.4), Inches(3.6), Inches(0.9), Inches(0.4))

# Chunks
chunk_x = Inches(4.5); chunk_y = Inches(2.2)
add_text(s, chunk_x, chunk_y, Inches(4), Inches(0.4),
         "Chunks (~300-500 words each)", size=14, bold=True, color=ACCENT)
for i in range(8):
    row = i // 2; col = i % 2
    x = chunk_x + Inches(col * 2.2)
    y = chunk_y + Inches(0.5 + row * 0.7)
    add_round(s, x, y, Inches(2.0), Inches(0.6), fill=ACCENT, line=ACCENT)
    add_text(s, x, y + Inches(0.12), Inches(2.0), Inches(0.4),
             f"Chunk {i+1}", size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

# Right side notes
add_round(s, Inches(9.4), Inches(2.2), Inches(3.4), Inches(3.5),
          fill=NAVY, line=NAVY)
add_text(s, Inches(9.6), Inches(2.35), Inches(3.0), Inches(0.4),
         "Good chunking:", size=14, bold=True, color=ACCENT)
add_bullets(s, Inches(9.6), Inches(2.85), Inches(3.0), Inches(2.8), [
    "Keeps related ideas together",
    "Small enough to be specific",
    "Big enough to have context",
    "Often with a little overlap",
], size=12, color=WHITE)

add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.5),
         "Think of it like cutting a textbook into flashcards — each card stands on its own.",
         size=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
footer(s, 7)


# ---------- Slide 8: Embeddings ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Step 2  —  Embeddings",
       "Turning meaning into numbers so a computer can compare it")

add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "An embedding is a list of numbers that represents the MEANING of text.",
         size=16, color=DARK_TEXT)

# Examples
examples = [
    ("“dog”",   "[0.21, -0.45, 0.88, …]"),
    ("“puppy”", "[0.23, -0.41, 0.86, …]"),
    ("“car”",   "[-0.72, 0.55, 0.10, …]"),
]
for i, (word, vec) in enumerate(examples):
    y = Inches(2.2 + i * 0.85)
    add_round(s, Inches(0.7), y, Inches(2.4), Inches(0.65),
              fill=TEAL, line=TEAL)
    add_text(s, Inches(0.7), y + Inches(0.13), Inches(2.4), Inches(0.4),
             word, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_round(s, Inches(3.3), y, Inches(5.5), Inches(0.65),
              fill=WHITE, line=GREY_TEXT)
    add_text(s, Inches(3.5), y + Inches(0.13), Inches(5.2), Inches(0.4),
             vec, size=14, color=DARK_TEXT)

add_text(s, Inches(0.7), Inches(4.9), Inches(12), Inches(0.5),
         "The magic:",
         size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(5.4), Inches(12), Inches(2), [
    "Words/phrases with SIMILAR MEANING get SIMILAR NUMBERS",
    "“dog” and “puppy” live close together; “car” lives far away",
    "So we can find related content even when the wording is different",
], size=14)

# Right: visualization
add_round(s, Inches(9.4), Inches(2.2), Inches(3.4), Inches(2.6),
          fill=WHITE, line=NAVY)
add_text(s, Inches(9.6), Inches(2.3), Inches(3.0), Inches(0.4),
         "Meaning-space", size=13, bold=True, color=NAVY)
# dots
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(3.0),
                        Inches(0.3), Inches(0.3))
dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT; dot.line.fill.background()
add_text(s, Inches(10.55), Inches(3.0), Inches(2), Inches(0.3),
         "dog", size=11, color=DARK_TEXT)
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.4), Inches(3.4),
                        Inches(0.3), Inches(0.3))
dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT; dot.line.fill.background()
add_text(s, Inches(10.75), Inches(3.4), Inches(2), Inches(0.3),
         "puppy", size=11, color=DARK_TEXT)
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.8), Inches(4.2),
                        Inches(0.3), Inches(0.3))
dot.fill.solid(); dot.fill.fore_color.rgb = TEAL; dot.line.fill.background()
add_text(s, Inches(11.2), Inches(4.2), Inches(0.6), Inches(0.3),
         "car", size=11, color=DARK_TEXT, align=PP_ALIGN.RIGHT)
footer(s, 8)


# ---------- Slide 9: Vector DB ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Step 3  —  Vector Database",
       "Where the numbered chunks live, and how we find the right ones fast")

add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "A vector database is built for ONE job:",
         size=16, color=DARK_TEXT)
add_text(s, Inches(0.7), Inches(1.9), Inches(12), Inches(0.6),
         "“Here's a vector. Find the closest ones in milliseconds.”",
         size=20, bold=True, color=ACCENT)

# Visual: query -> DB -> top-k
add_round(s, Inches(0.7), Inches(3.0), Inches(2.6), Inches(1.5),
          fill=TEAL, line=TEAL)
add_text(s, Inches(0.7), Inches(3.2), Inches(2.6), Inches(0.4),
         "Question", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(3.65), Inches(2.6), Inches(0.7),
         "“What's the\nleave policy?”",
         size=13, color=WHITE, align=PP_ALIGN.CENTER)

add_arrow(s, Inches(3.4), Inches(3.6), Inches(0.7), Inches(0.3))

add_round(s, Inches(4.3), Inches(3.0), Inches(3.5), Inches(1.5),
          fill=NAVY, line=NAVY)
add_text(s, Inches(4.3), Inches(3.15), Inches(3.5), Inches(0.4),
         "Vector DB", size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.3), Inches(3.55), Inches(3.5), Inches(0.4),
         "(thousands of chunks)", size=11, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.3), Inches(3.95), Inches(3.5), Inches(0.4),
         "Pinecone · Chroma · FAISS · Weaviate", size=10, color=ACCENT,
         align=PP_ALIGN.CENTER)

add_arrow(s, Inches(7.9), Inches(3.6), Inches(0.7), Inches(0.3))

add_round(s, Inches(8.8), Inches(3.0), Inches(4.0), Inches(1.5),
          fill=ACCENT, line=ACCENT)
add_text(s, Inches(8.8), Inches(3.15), Inches(4.0), Inches(0.4),
         "Top 3-5 chunks", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.8), Inches(3.6), Inches(4.0), Inches(0.7),
         "The MOST RELEVANT\npieces of your documents",
         size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_bullets(s, Inches(0.7), Inches(5.0), Inches(12), Inches(2), [
    "No keyword matching — this is meaning matching",
    "Works even if the user asks the question differently than the doc was written",
    "Returns the top K (usually 3-10) chunks with similarity scores",
], size=14)
footer(s, 9)


# ---------- Slide 10: Putting it together ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Putting It All Together  —  One Query",
       "Follow the path of a single question through the system")

# Big horizontal flow
nodes = [
    ("User\nquestion",    TEAL),
    ("Embed\nquestion",   NAVY),
    ("Search\nVector DB", ACCENT),
    ("Top chunks\nretrieved", NAVY),
    ("Stuff into\nprompt", TEAL),
    ("LLM\nanswer", ACCENT),
]
n = len(nodes)
node_w = Inches(1.8); gap = Inches(0.15)
total = node_w * n + gap * (n - 1)
sx = (SW - total) / 2; y = Inches(2.0)
for i, (t, col) in enumerate(nodes):
    x = sx + i * (node_w + gap)
    add_round(s, x, y, node_w, Inches(1.6), fill=col, line=col)
    add_text(s, x + Inches(0.1), y + Inches(0.4), node_w - Inches(0.2), Inches(1.0),
             t, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < n - 1:
        add_arrow(s, x + node_w - Inches(0.05), y + Inches(0.65),
                  gap + Inches(0.1), Inches(0.25), color=GREY_TEXT)

# Final prompt example
add_text(s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.4),
         "The prompt the LLM actually sees:",
         size=14, bold=True, color=NAVY)
add_round(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(2.3),
          fill=WHITE, line=NAVY)
add_text(s, Inches(0.9), Inches(4.75), Inches(11.5), Inches(2.1),
         "You are a helpful assistant. Use ONLY the context below to answer.\n\n"
         "CONTEXT:\n"
         "  • Chunk 1: “Faculty are entitled to 12 days of casual leave…”\n"
         "  • Chunk 2: “Leave applications must be submitted via the HR portal…”\n"
         "  • Chunk 3: “Medical leave requires a certificate after 3 days…”\n\n"
         "QUESTION: What's the leave policy for faculty?",
         size=12, color=DARK_TEXT)
footer(s, 10)


# ---------- Slide 11: Light technical walk - code snippet ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "What the Code Roughly Looks Like",
       "You don't have to write it — but here's how short it is")

add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.4),
         "Using LangChain + OpenAI + Chroma (~15 lines):",
         size=14, bold=True, color=NAVY)

code = (
    "from langchain.document_loaders import PyPDFLoader\n"
    "from langchain.text_splitter import RecursiveCharacterTextSplitter\n"
    "from langchain.embeddings import OpenAIEmbeddings\n"
    "from langchain.vectorstores import Chroma\n"
    "from langchain.chains import RetrievalQA\n"
    "from langchain.chat_models import ChatOpenAI\n\n"
    "# 1. Load + chunk\n"
    "docs = PyPDFLoader('policy.pdf').load()\n"
    "chunks = RecursiveCharacterTextSplitter(chunk_size=500).split_documents(docs)\n\n"
    "# 2. Embed + store\n"
    "db = Chroma.from_documents(chunks, OpenAIEmbeddings())\n\n"
    "# 3. Ask\n"
    "qa = RetrievalQA.from_chain_type(ChatOpenAI(), retriever=db.as_retriever())\n"
    "print(qa.run('What's the leave policy?'))"
)
add_round(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(4.6),
          fill=NAVY, line=NAVY)
tb = add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.4),
              code, size=13, color=WHITE)
# Change font to monospace
for p in tb.text_frame.paragraphs:
    for r in p.runs:
        r.font.name = "Consolas"

add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
         "That's it. Six imports, three steps, one production-grade pattern.",
         size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
footer(s, 11)


# ---------- Slide 11.5: LIVE DEMO — PolicyBot ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "🎬  LIVE DEMO  —  PolicyBot",
       "Let's run RAG on a real document, right now")

# Left: what we'll do
add_round(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.4),
          fill=WHITE, line=ACCENT)
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.7), fill=ACCENT)
add_text(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.5),
         "What we'll do",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
steps = [
    "📄  Upload a Faculty Handbook PDF",
    "✂️  Watch it get chunked & embedded",
    "❓  Ask questions in plain English",
    "✅  See the answer",
    "📑  See the exact chunks the AI used",
]
add_bullets(s, Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.0),
            steps, size=15)
add_text(s, Inches(0.8), Inches(6.0), Inches(5.5), Inches(0.5),
         "policybot.py  •  ~80 lines  •  same Streamlit feel",
         size=12, color=GREY_TEXT, bold=True)

# Right: sample questions to try
add_round(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.4),
          fill=NAVY, line=NAVY)
add_text(s, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.5),
         "Try these questions:",
         size=18, bold=True, color=ACCENT)
questions = [
    "“How many casual leaves do faculty get?”",
    "“What's the work-from-home policy?”",
    "“What's required to become Associate Professor?”",
    "“Can faculty use ChatGPT for student data?”",
    "“What happens if I'm late 3 times in a month?”",
]
add_bullets(s, Inches(7.0), Inches(2.4), Inches(5.6), Inches(4.0),
            questions, size=13, color=WHITE, bullet="›")
add_text(s, Inches(7.0), Inches(6.0), Inches(5.6), Inches(0.5),
         "👀  Watch the “Retrieved Chunks” panel — that's the magic.",
         size=12, color=ACCENT, bold=True)
footer(s, 12)


# ---------- Slide 12: When to use what ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Prompt-Stuffing vs RAG vs Fine-Tuning",
       "Three tools. Pick the right one.")

cols = [
    ("Prompt-Stuffing", TEAL,
     "Paste data\ninto the prompt",
     ["Tiny datasets", "Quick demos", "Our morning app"],
     "❌ Doesn't scale"),
    ("RAG", ACCENT,
     "Retrieve, then\ngenerate",
     ["Large doc sets", "Always up-to-date", "Need citations", "Most enterprise AI"],
     "✅ The sweet spot"),
    ("Fine-Tuning", NAVY,
     "Retrain the\nmodel itself",
     ["Custom style/tone", "Domain expertise", "Rare — expensive"],
     "⚠️ Heavy lift"),
]
for i, (name, col, what, when, verdict) in enumerate(cols):
    x = Inches(0.5 + i * 4.27)
    y = Inches(1.4)
    add_round(s, x, y, Inches(4.05), Inches(5.5), fill=WHITE, line=col)
    add_rect(s, x, y, Inches(4.05), Inches(0.8), fill=col)
    add_text(s, x, y + Inches(0.15), Inches(4.05), Inches(0.5),
             name, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(1.0), Inches(3.7), Inches(0.9),
             what, size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + Inches(2.1), Inches(3.4), Inches(0.4),
             "Best when:", size=13, bold=True, color=NAVY)
    add_bullets(s, x + Inches(0.3), y + Inches(2.5), Inches(3.5), Inches(2.3),
                when, size=12)
    add_text(s, x + Inches(0.3), y + Inches(4.8), Inches(3.5), Inches(0.5),
             verdict, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
footer(s, 13)


# ---------- Slide 14: Real-world faculty examples ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "What Could YOU Build with RAG?",
       "Real ideas for this campus")

ideas = [
    ("📋", "Policy Bot",
     "Ask: “What's the leave rule for visiting faculty?”\nSearches HR policy PDFs and answers with citations."),
    ("📚", "Syllabus Q&A",
     "Students ask any question — bot answers strictly\nfrom course materials you uploaded."),
    ("🎓", "NAAC / NBA Helper",
     "Searches years of accreditation docs to draft\nresponses to criteria questions."),
    ("🔬", "Research Companion",
     "Upload 50 papers. Ask cross-paper questions.\nGet summaries with sources."),
]
for i, (ic, t, d) in enumerate(ideas):
    row = i // 2; col = i % 2
    x = Inches(0.5 + col * 6.27)
    y = Inches(1.5 + row * 2.6)
    add_round(s, x, y, Inches(6.0), Inches(2.4), fill=WHITE, line=ACCENT)
    add_text(s, x + Inches(0.25), y + Inches(0.3), Inches(1.0), Inches(0.8),
             ic, size=36, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.3), y + Inches(0.35), Inches(4.5), Inches(0.5),
             t, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(1.3), y + Inches(0.95), Inches(4.5), Inches(1.4),
             d, size=12, color=DARK_TEXT)
footer(s, 14)


# ---------- Slide 15: Gotchas ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Honest Caveats",
       "RAG is powerful — but it's not magic")

caveats = [
    ("📐", "Chunk size matters",
     "Too small loses context. Too big drowns the LLM. It's an art."),
    ("🔍", "Retrieval can miss",
     "If your question uses very different language from the doc, top chunks may not be the best ones."),
    ("📑", "Garbage in, garbage out",
     "If your source documents are messy, scanned, or outdated — so is the answer."),
    ("🔒", "Privacy still applies",
     "Sending content to an API = data leaves your campus. Use local models for sensitive data."),
]
for i, (ic, t, d) in enumerate(caveats):
    y = Inches(1.4 + i * 1.35)
    add_round(s, Inches(0.6), y, Inches(12.1), Inches(1.2),
              fill=WHITE, line=TEAL)
    add_text(s, Inches(0.85), y + Inches(0.25), Inches(0.8), Inches(0.8),
             ic, size=32, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.9), y + Inches(0.15), Inches(10.5), Inches(0.5),
             t, size=18, bold=True, color=NAVY)
    add_text(s, Inches(1.9), y + Inches(0.6), Inches(10.5), Inches(0.6),
             d, size=13, color=DARK_TEXT)
footer(s, 15)


# ---------- Slide 16: Hallucination Shock Demo ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "☠️  Why RAG Exists  —  The Hallucination Problem",
       "Ask a plain LLM about YOUR world — watch it lie confidently")

# Left: closed-book LLM
add_round(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.4),
          fill=WHITE, line=RED_NO)
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.7), fill=RED_NO)
add_text(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.5),
         "Plain LLM  —  No RAG", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.5),
         "You ask:", size=13, bold=True, color=GREY_TEXT)
add_text(s, Inches(0.8), Inches(2.7), Inches(5.5), Inches(0.6),
         "“How many casual leaves do\nHITS faculty get?”",
         size=15, bold=True, color=NAVY)
add_text(s, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.5),
         "It answers (confidently):", size=13, bold=True, color=GREY_TEXT)
add_round(s, Inches(0.8), Inches(4.4), Inches(5.4), Inches(2.1),
          fill=SOFT_BLUE, line=GREY_TEXT)
add_text(s, Inches(1.0), Inches(4.55), Inches(5.0), Inches(1.9),
         "“Faculty get 15 casual leaves\nper year as per UGC norms…”\n\n⚠️  Sounds right.\n⚠️  Cites nothing.\n⚠️  Completely made up.",
         size=13, color=DARK_TEXT)

# Right: RAG
add_round(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.4),
          fill=WHITE, line=GREEN_OK)
add_rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.7), fill=GREEN_OK)
add_text(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.5),
         "LLM + RAG", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.5), Inches(0.5),
         "Same question — grounded answer:", size=13, bold=True, color=GREY_TEXT)
add_round(s, Inches(7.1), Inches(2.8), Inches(5.4), Inches(3.6),
          fill=SOFT_BLUE, line=GREEN_OK)
add_text(s, Inches(7.3), Inches(2.95), Inches(5.0), Inches(3.4),
         "“As per the HITS Faculty\nHandbook 2024 (Sec 4.2),\nfull-time faculty are entitled\nto 12 days of casual leave.”\n\n✅  Pulled from real doc\n✅  Cites section\n✅  Verifiable",
         size=13, color=DARK_TEXT)

add_text(s, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4),
         "Hallucination isn’t a bug — it’s the default. RAG fixes it by giving the LLM real facts to ground on.",
         size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
footer(s, 16)


# ---------- Slide 17: Chunk Size Matters ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "The Chunk-Size Sweet Spot",
       "Same document, three different chunk sizes — very different results")

cols2 = [
    ("Too Small", "~100 chars", RED_NO,
     ["Loses surrounding context",
      "“leave” alone—can’t tell which kind",
      "Many tiny chunks → fragmented answers"],
     "❌  Specific but meaningless"),
    ("Just Right", "~500 chars", GREEN_OK,
     ["One idea per chunk",
      "Enough context to make sense",
      "Small enough to be precise",
      "Often with 10-20% overlap"],
     "✅  The sweet spot for most docs"),
    ("Too Big", "~3000 chars", RED_NO,
     ["Mixes multiple topics together",
      "LLM gets distracted by noise",
      "Wastes tokens → more cost",
      "Less precise retrieval"],
     "❌  Drowns the signal"),
]
for i, (name, sz, col, pts, verdict) in enumerate(cols2):
    x = Inches(0.5 + i * 4.27)
    y = Inches(1.4)
    add_round(s, x, y, Inches(4.05), Inches(5.4), fill=WHITE, line=col)
    add_rect(s, x, y, Inches(4.05), Inches(0.9), fill=col)
    add_text(s, x, y + Inches(0.1), Inches(4.05), Inches(0.45),
             name, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.5), Inches(4.05), Inches(0.4),
             sz, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullets(s, x + Inches(0.3), y + Inches(1.2), Inches(3.5), Inches(3.4),
                pts, size=12)
    add_text(s, x + Inches(0.2), y + Inches(4.7), Inches(3.7), Inches(0.5),
             verdict, size=13, bold=True, color=col, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4),
         "Rule of thumb: start with 500 chars + 50 overlap. Then tune by trying real questions.",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
footer(s, 17)


# ---------- Slide 18: Advanced RAG ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Going Further  —  Advanced RAG Patterns",
       "Tricks that production systems add on top of the basics")

adv = [
    ("🔀", "Hybrid Search",
     "Combine keyword (BM25) + vector search.\nKeyword catches exact terms like ‘Section 4.2’; vectors catch meaning."),
    ("🎯", "Re-Ranking",
     "Retrieve top 50, then a smarter (slower) model re-sorts them to find the BEST 5.\nHuge accuracy boost."),
    ("📄", "Parent-Document",
     "Search small chunks for precision, but send the LARGER surrounding paragraph to the LLM for context."),
    ("🧠", "Query Rewriting",
     "LLM rephrases the user’s vague question into a better search query before retrieving."),
    ("🗂️", "Metadata Filtering",
     "Tag chunks with course, year, department. Filter BEFORE searching: ‘only CS dept docs from 2024’."),
    ("🔁", "Multi-Hop / Agentic RAG",
     "For complex questions, the system retrieves, reasons, then retrieves AGAIN. (Bridge to Day 2.)"),
]
for i, (ic, t, d) in enumerate(adv):
    row = i // 2; col = i % 2
    x = Inches(0.5 + col * 6.27)
    y = Inches(1.4 + row * 1.85)
    add_round(s, x, y, Inches(6.0), Inches(1.7), fill=WHITE, line=PURPLE)
    add_text(s, x + Inches(0.2), y + Inches(0.3), Inches(0.9), Inches(0.8),
             ic, size=30, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.2), y + Inches(0.2), Inches(4.7), Inches(0.5),
             t, size=16, bold=True, color=NAVY)
    add_text(s, x + Inches(1.2), y + Inches(0.7), Inches(4.7), Inches(1.0),
             d, size=11, color=DARK_TEXT)
footer(s, 18)


# ---------- Slide 19: How do we know RAG works? ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "How Do You Know Your RAG Actually Works?",
       "Evaluation — the part everyone skips, until it bites them")

add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "A RAG system can fail at two places. Test BOTH.",
         size=16, bold=True, color=NAVY)

# Two pillars
pillars = [
    ("🔍", "Retrieval Quality",
     "Did we fetch the RIGHT chunks?", TEAL,
     ["Hit Rate  —  is the correct chunk in top-K?",
      "MRR  —  how high did the correct chunk rank?",
      "Recall@K  —  fraction of relevant chunks found"]),
    ("✍️", "Answer Quality",
     "Did the LLM use them well?", ACCENT,
     ["Faithfulness  —  is the answer grounded in the chunks?",
      "Relevance  —  does it actually answer the question?",
      "Citation Accuracy  —  are the sources correct?"]),
]
for i, (ic, t, sub, col, pts) in enumerate(pillars):
    x = Inches(0.5 + i * 6.27)
    y = Inches(2.1)
    add_round(s, x, y, Inches(6.0), Inches(4.4), fill=WHITE, line=col)
    add_rect(s, x, y, Inches(6.0), Inches(0.9), fill=col)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(0.8), Inches(0.6),
             ic, size=28, color=WHITE)
    add_text(s, x + Inches(1.1), y + Inches(0.1), Inches(4.7), Inches(0.45),
             t, size=18, bold=True, color=WHITE)
    add_text(s, x + Inches(1.1), y + Inches(0.5), Inches(4.7), Inches(0.4),
             sub, size=12, color=WHITE)
    add_bullets(s, x + Inches(0.3), y + Inches(1.2), Inches(5.5), Inches(3.0),
                pts, size=13)

add_text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
         "Tools: RAGAS  ·  LangSmith  ·  TruLens  —  or build a small ‘golden Q&A’ set by hand.",
         size=13, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
footer(s, 19)


# ---------- Slide 20: Privacy & Local RAG ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Privacy-First RAG  —  Runs On YOUR Laptop",
       "For student data, exam papers, internal HR — nothing should leave campus")

# Left: the worry
add_round(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.0),
          fill=WHITE, line=RED_NO)
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.7), fill=RED_NO)
add_text(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.5),
         "The Worry  😨", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_bullets(s, Inches(0.8), Inches(2.3), Inches(5.4), Inches(4.0),
            ["Sending student answers to OpenAI?",
             "Uploading internal HR policies to a US server?",
             "What about NAAC confidentiality?",
             "Data leaves campus = compliance nightmare."], size=14)

# Right: the fix
add_round(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.0),
          fill=WHITE, line=GREEN_OK)
add_rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.7), fill=GREEN_OK)
add_text(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.5),
         "The Fix  —  Local Stack  🔒", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
local_stack = [
    ("LLM", "Ollama  ·  Llama 3  ·  Mistral"),
    ("Embeddings", "BGE  ·  all-MiniLM (open source)"),
    ("Vector DB", "Chroma  ·  FAISS  (on disk)"),
    ("Network", "Zero internet calls required"),
]
for i, (k, v) in enumerate(local_stack):
    y = Inches(2.3 + i * 0.85)
    add_round(s, Inches(7.0), y, Inches(2.0), Inches(0.65),
              fill=NAVY, line=NAVY)
    add_text(s, Inches(7.0), y + Inches(0.15), Inches(2.0), Inches(0.4),
             k, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.2), y + Inches(0.2), Inches(3.5), Inches(0.4),
             v, size=12, color=DARK_TEXT)

add_text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
         "Same RAG architecture —  just swap each piece for an open-source equivalent.",
         size=14, bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)
footer(s, 20)


# ---------- Slide 21: Your Turn ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "🚀  Your Turn  —  Build One This Week",
       "A simple 5-step challenge for your own course or department")

challenge = [
    ("1", "Pick ONE document",
     "A syllabus, a policy doc, your lecture notes — anything you know well."),
    ("2", "List 5 real questions",
     "Questions students or colleagues genuinely ask you about that doc."),
    ("3", "Drop it into PolicyBot",
     "Use today’s demo code. No new setup needed. Index your PDF in 1 minute."),
    ("4", "Ask your 5 questions",
     "Note where it shines  ✅  and where it stumbles  ⚠️."),
    ("5", "Share back tomorrow",
     "Bring your finding to Day 2 — we’ll turn the best ones into agents."),
]
for i, (n, t, d) in enumerate(challenge):
    y = Inches(1.4 + i * 1.05)
    add_round(s, Inches(0.6), y, Inches(12.1), Inches(0.95),
              fill=WHITE, line=ACCENT)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.2),
                              Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    add_text(s, Inches(0.85), y + Inches(0.23), Inches(0.55), Inches(0.5),
             n, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.7), y + Inches(0.12), Inches(11.0), Inches(0.45),
             t, size=16, bold=True, color=NAVY)
    add_text(s, Inches(1.7), y + Inches(0.52), Inches(11.0), Inches(0.45),
             d, size=12, color=DARK_TEXT)

add_text(s, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4),
         "The faculty who actually TRY it tonight will get the most out of Day 2.",
         size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
footer(s, 21)


# ---------- Slide 22: Recap ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Recap  —  What You Now Know About RAG")
points = [
    ("1", "RAG = Retrieve first, then Generate"),
    ("2", "It solves the “too much data, too few tokens” problem"),
    ("3", "Documents → chunks → embeddings → vector DB → retrieval → LLM"),
    ("4", "Embeddings let us find by MEANING, not keywords"),
    ("5", "RAG is the default architecture for enterprise AI today"),
    ("6", "Tomorrow we'll see what happens when AI starts to ACT — Agentic AI"),
]
for i, (n, t) in enumerate(points):
    y = Inches(1.4 + i * 0.85)
    add_round(s, Inches(0.6), y, Inches(12.1), Inches(0.75),
              fill=WHITE, line=TEAL)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.1),
                              Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    add_text(s, Inches(0.8), y + Inches(0.13), Inches(0.55), Inches(0.5),
             n, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.55), y + Inches(0.17), Inches(11.0), Inches(0.5),
             t, size=15, color=DARK_TEXT)
footer(s, 22)


# ---------- Slide 23: Bridge to Day 2 ----------
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_rect(s, 0, Inches(3.2), SW, Inches(0.08), fill=ACCENT)
add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.5),
         "TOMORROW  —  DAY 2",
         size=16, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1.4),
         "AI that thinks  is useful.",
         size=38, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(1.4),
         "AI that acts  is transformative.",
         size=38, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.6),
         "Day 2  —  Agentic AI",
         size=28, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.6),
         "Tools • Function-calling • Autonomous agents • Multi-agent systems",
         size=16, color=RGBColor(0xCF, 0xE3, 0xF5))
add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
         "Thank you  —  See you tomorrow!",
         size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


import os
out = "Day1_Part3_RAG.pptx"
try:
    prs.save(out)
except PermissionError:
    out = "Day1_Part3_RAG_v2.pptx"
    prs.save(out)
print(f"Saved: {out}  ({os.path.abspath(out)})")
