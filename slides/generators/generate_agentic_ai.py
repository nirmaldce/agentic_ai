"""
Generates 'Day2_Agentic_AI.pptx' — Day 2 of the workshop.
Run: py generate_agentic_ai.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY      = RGBColor(0x0B, 0x2A, 0x4A)
TEAL      = RGBColor(0x00, 0x9B, 0x9E)
ACCENT    = RGBColor(0xF4, 0xA2, 0x61)
PURPLE    = RGBColor(0x7A, 0x4E, 0xAB)
MAGENTA   = RGBColor(0xD1, 0x3F, 0x7C)
LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x2D, 0x3D)
GREY_TEXT = RGBColor(0x55, 0x66, 0x77)
SOFT_BLUE = RGBColor(0xE8, 0xF1, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=LIGHT_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    return bg

def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    return s

def add_round(slide, x, y, w, h, fill=WHITE, line=NAVY, adj=0.12):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = adj
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(1.25)
    return s

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def add_bullets(slide, x, y, w, h, items, size=16, color=DARK_TEXT, bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run(); r.text = f"{bullet}  {item}"
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def add_arrow(slide, x, y, w, h, color=TEAL):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a

def add_circ_arrow(slide, x, y, w, h, color=PURPLE):
    a = slide.shapes.add_shape(MSO_SHAPE.CIRCULAR_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a

def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(1.0), fill=NAVY)
    add_rect(slide, 0, Inches(1.0), SW, Inches(0.08), fill=PURPLE)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.62), Inches(12), Inches(0.4),
                 subtitle, size=14, color=RGBColor(0xCF, 0xE3, 0xF5))

def footer(slide, page):
    add_text(slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.3),
             "Day 2 — Agentic AI", size=10, color=GREY_TEXT)
    add_text(slide, Inches(12.3), Inches(7.15), Inches(1), Inches(0.3),
             str(page), size=10, color=GREY_TEXT, align=PP_ALIGN.RIGHT)


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_rect(s, 0, Inches(3.2), SW, Inches(1.2), fill=PURPLE)
add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "DAY 2  •  THE WOW DAY", size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(1.4),
         "Agentic AI", size=72, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(2.7), Inches(12), Inches(0.7),
         "From AI that answers  →  to AI that ACTS",
         size=28, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.5), Inches(12), Inches(0.7),
         "Tools. Agents. Multi-agent teams. Live demo.",
         size=20, color=RGBColor(0xCF, 0xE3, 0xF5))
add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
         "Two-Day Workshop on Agentic AI",
         size=14, color=ACCENT, bold=True)


# ---------- Slide 2: Recap ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Yesterday, We Built AI That ANSWERS",
       "Today, we build AI that ACTS")

# Yesterday timeline
add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "Day 1 recap:", size=18, bold=True, color=NAVY)
day1 = [
    ("Session 1", "Intro to AI",     "We understood how AI thinks"),
    ("Session 2", "Faculty Asst",    "Excel → Prompt → LLM answers"),
    ("Session 3", "RAG / PolicyBot", "Big docs → retrieve → grounded answers"),
]
for i, (s1, s2, s3) in enumerate(day1):
    y = Inches(2.1 + i * 0.85)
    add_round(s, Inches(0.7), y, Inches(12), Inches(0.7), fill=WHITE, line=TEAL)
    add_text(s, Inches(0.9), y + Inches(0.18), Inches(2.0), Inches(0.4),
             s1, size=14, bold=True, color=TEAL)
    add_text(s, Inches(2.9), y + Inches(0.15), Inches(3.5), Inches(0.4),
             s2, size=15, bold=True, color=NAVY)
    add_text(s, Inches(6.6), y + Inches(0.17), Inches(6.0), Inches(0.4),
             s3, size=13, color=GREY_TEXT)

# Connector
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5),
         "But every single time…",
         size=18, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.7),
         "YOU had to take the next action.",
         size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.5),
         "What if the AI took the action itself?",
         size=20, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
footer(s, 2)


# ---------- Slide 3: Today's promise ----------
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_rect(s, 0, Inches(3.5), SW, Inches(0.08), fill=ACCENT)
add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.7),
         "By the end of today, you will see",
         size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(1.5),
         "an AI team",
         size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.7),
         "look at your students,",
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.8), Inches(12), Inches(0.7),
         "find the ones who need help,",
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.7),
         "and do the whole follow-up — by itself.",
         size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ---------- Slide 4: What is an Agent ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "So… What is an Agent?",
       "It's an LLM that's been given a job, tools, and permission to use them")

# Formula
add_round(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(1.0),
          fill=NAVY, line=NAVY)
add_text(s, Inches(0.7), Inches(1.65), Inches(11.9), Inches(0.6),
         "Agent  =  LLM  +  Tools  +  Memory  +  Goal",
         size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Four parts
parts = [
    ("🧠", "LLM",    "The brain — reasons & decides",      TEAL),
    ("🛠️", "Tools",  "Hands — APIs, code, search, email",  ACCENT),
    ("📓", "Memory", "Remembers what it did & saw",         PURPLE),
    ("🎯", "Goal",   "The job you gave it",                NAVY),
]
for i, (ic, name, desc, col) in enumerate(parts):
    x = Inches(0.7 + i * 3.1)
    y = Inches(2.8)
    add_round(s, x, y, Inches(2.9), Inches(3.9), fill=WHITE, line=col)
    add_rect(s, x, y, Inches(2.9), Inches(0.7), fill=col)
    add_text(s, x, y + Inches(0.13), Inches(2.9), Inches(0.5),
             name, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.0), Inches(2.9), Inches(1.0),
             ic, size=56, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(2.5), Inches(2.5), Inches(1.4),
             desc, size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)
footer(s, 4)


# ---------- Slide 5: Assistant vs Agent ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Assistant  vs.  Agent",
       "Same brain. Very different behavior.")

# Left: Assistant
add_round(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.4),
          fill=WHITE, line=TEAL)
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.8), fill=TEAL)
add_text(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.5),
         "🤖  ASSISTANT  (yesterday)", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_bullets(s, Inches(0.8), Inches(2.5), Inches(5.6), Inches(4.0), [
    "Waits for your question",
    "Returns ONE answer",
    "Stops there",
    "You decide what to do next",
    "Reactive",
], size=15)

# Right: Agent
add_round(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.4),
          fill=WHITE, line=ACCENT)
add_rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.8), fill=ACCENT)
add_text(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.5),
         "🚀  AGENT  (today)", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_bullets(s, Inches(7.1), Inches(2.5), Inches(5.6), Inches(4.0), [
    "You give it a GOAL, not a question",
    "Plans its own steps",
    "Uses tools (search, email, calendar, code)",
    "Adapts based on what it sees",
    "Keeps going until the goal is done",
    "Proactive",
], size=15)
footer(s, 5)


# ---------- Slide 6: The ReAct loop ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "How Agents Think  —  The ReAct Loop",
       "Reason  →  Act  →  Observe  →  Repeat")

# 4 circles in a cycle
nodes = [
    ("THINK",   "What should I do next?",   TEAL,   Inches(3.5), Inches(2.0)),
    ("ACT",     "Use a tool",               ACCENT, Inches(8.0), Inches(2.0)),
    ("OBSERVE", "What did I get back?",     PURPLE, Inches(8.0), Inches(4.8)),
    ("REPEAT",  "Goal done? If not, loop", NAVY,   Inches(3.5), Inches(4.8)),
]
for label, desc, col, x, y in nodes:
    add_round(s, x, y, Inches(2.0), Inches(2.0), fill=col, line=col, adj=0.5)
    add_text(s, x, y + Inches(0.55), Inches(2.0), Inches(0.5),
             label, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(1.05), Inches(1.7), Inches(0.9),
             desc, size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Arrows connecting them clockwise
add_arrow(s, Inches(5.6), Inches(2.7), Inches(2.3), Inches(0.5), color=GREY_TEXT)
arr2 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(8.5), Inches(4.1),
                           Inches(1.0), Inches(0.6))
arr2.fill.solid(); arr2.fill.fore_color.rgb = GREY_TEXT
arr2.line.fill.background()
arr3 = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(5.6), Inches(5.5),
                           Inches(2.3), Inches(0.5))
arr3.fill.solid(); arr3.fill.fore_color.rgb = GREY_TEXT
arr3.line.fill.background()
arr4 = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(4.0), Inches(4.1),
                           Inches(1.0), Inches(0.6))
arr4.fill.solid(); arr4.fill.fore_color.rgb = GREY_TEXT
arr4.line.fill.background()

# Tagline
add_text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.4),
         "Every modern AI agent — from ChatGPT plugins to Cursor to Devin — runs this loop.",
         size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
footer(s, 6)


# ---------- Slide 7: Example trace ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "What an Agent's 'Thoughts' Look Like",
       "A real ReAct trace — readable, debuggable, transparent")

trace = [
    ("GOAL",    "Identify weak students and email their parents.", NAVY),
    ("THINK",   "First I need the student data.",                  TEAL),
    ("ACT",     "tool: read_excel('students.xlsx')",               ACCENT),
    ("OBSERVE", "Got 47 rows. Columns: Name, Marks, Attendance.",  PURPLE),
    ("THINK",   "I'll flag students with marks<40 OR attendance<75.", TEAL),
    ("ACT",     "tool: filter(marks<40 or attendance<75)",         ACCENT),
    ("OBSERVE", "5 students flagged.",                             PURPLE),
    ("THINK",   "Now I'll draft an email per student.",            TEAL),
    ("ACT",     "tool: draft_email(student=…) ×5",                 ACCENT),
    ("OBSERVE", "5 drafts ready. Goal achieved.",                  PURPLE),
]
y0 = Inches(1.35)
row_h = Inches(0.5)
for i, (tag, content, col) in enumerate(trace):
    y = y0 + i * row_h
    add_round(s, Inches(0.7), y + Inches(0.05), Inches(1.4), Inches(0.4),
              fill=col, line=col, adj=0.3)
    add_text(s, Inches(0.7), y + Inches(0.08), Inches(1.4), Inches(0.35),
             tag, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.25), y + Inches(0.08), Inches(10.5), Inches(0.4),
             content, size=13, color=DARK_TEXT)
footer(s, 7)


# ---------- Slide 8: Tools / Function calling ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Tools  —  The Agent's Superpower",
       "Anything you can call from code, an agent can call too")

# Visual: LLM in middle, tools around
add_round(s, Inches(5.5), Inches(3.0), Inches(2.5), Inches(1.5),
          fill=NAVY, line=NAVY)
add_text(s, Inches(5.5), Inches(3.3), Inches(2.5), Inches(0.5),
         "🧠  LLM", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(5.5), Inches(3.8), Inches(2.5), Inches(0.4),
         "(the brain)", size=12, color=ACCENT, align=PP_ALIGN.CENTER)

tools = [
    ("📊", "Read Excel",   Inches(1.0),  Inches(1.5)),
    ("🌐", "Web search",   Inches(10.0), Inches(1.5)),
    ("✉️", "Send email",   Inches(1.0),  Inches(3.0)),
    ("📅", "Calendar",     Inches(10.0), Inches(3.0)),
    ("🗄️", "Database",    Inches(1.0),  Inches(4.5)),
    ("🐍", "Run code",     Inches(10.0), Inches(4.5)),
    ("📄", "Read PDF",     Inches(1.0),  Inches(6.0)),
    ("💬", "Slack/Teams",  Inches(10.0), Inches(6.0)),
]
for ic, name, x, y in tools:
    add_round(s, x, y, Inches(2.3), Inches(0.8), fill=WHITE, line=ACCENT)
    add_text(s, x + Inches(0.1), y + Inches(0.15), Inches(0.6), Inches(0.5),
             ic, size=20, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.7), y + Inches(0.22), Inches(1.5), Inches(0.4),
             name, size=13, bold=True, color=NAVY)

add_text(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
         "The LLM decides WHICH tool to call, with WHAT arguments. We just register them.",
         size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
footer(s, 8)


# ---------- Slide 9: From one agent to many ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "One Agent is Cool. A TEAM of Agents is Magic.",
       "Multi-agent systems — like a team of specialists")

add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "Why split the job across multiple agents?",
         size=18, bold=True, color=NAVY)

reasons = [
    ("🎯", "Specialization",
     "Each agent has ONE clear job & prompt. Better results than one mega-prompt."),
    ("🔁", "Reusability",
     "An 'Email Drafter' agent works for students, parents, vendors — anywhere."),
    ("🧩", "Composability",
     "Mix and match agents for new workflows without rewriting."),
    ("🔍", "Transparency",
     "When something goes wrong, you can see exactly WHICH agent failed."),
]
for i, (ic, t, d) in enumerate(reasons):
    row = i // 2; col = i % 2
    x = Inches(0.6 + col * 6.3)
    y = Inches(2.2 + row * 2.2)
    add_round(s, x, y, Inches(6.0), Inches(2.0), fill=WHITE, line=PURPLE)
    add_text(s, x + Inches(0.25), y + Inches(0.3), Inches(1.0), Inches(0.8),
             ic, size=36, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.4), y + Inches(0.3), Inches(4.4), Inches(0.5),
             t, size=18, bold=True, color=PURPLE)
    add_text(s, x + Inches(1.4), y + Inches(0.85), Inches(4.4), Inches(1.0),
             d, size=12, color=DARK_TEXT)
footer(s, 9)


# ---------- Slide 10: Multi-agent patterns ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "How Agents Work Together  —  3 Patterns",
       "Sequential. Hierarchical. Debate.")

patterns = [
    ("Sequential",
     "A pipeline — each agent passes output to the next.",
     "Analyst → Writer → Editor",
     TEAL),
    ("Hierarchical",
     "A manager agent delegates to specialist workers.",
     "Manager → (Researcher, Coder, Writer)",
     ACCENT),
    ("Debate",
     "Agents argue / critique each other to improve quality.",
     "Proposer ⇄ Critic ⇄ Judge",
     PURPLE),
]
for i, (name, desc, ex, col) in enumerate(patterns):
    x = Inches(0.5 + i * 4.27)
    y = Inches(1.4)
    add_round(s, x, y, Inches(4.05), Inches(5.5), fill=WHITE, line=col)
    add_rect(s, x, y, Inches(4.05), Inches(0.8), fill=col)
    add_text(s, x, y + Inches(0.15), Inches(4.05), Inches(0.5),
             name, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + Inches(1.1), Inches(3.5), Inches(1.5),
             desc, size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    add_round(s, x + Inches(0.3), y + Inches(3.0), Inches(3.5), Inches(1.3),
              fill=LIGHT_BG, line=col)
    add_text(s, x + Inches(0.3), y + Inches(3.25), Inches(3.5), Inches(0.4),
             "Example", size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + Inches(3.65), Inches(3.5), Inches(0.8),
             ex, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + Inches(4.6), Inches(3.5), Inches(0.5),
             "✨ Our demo uses this →" if i == 0 else "",
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
footer(s, 10)


# ---------- Slide 11: The demo team ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Meet the Team  —  Faculty Operations Agents",
       "4 specialists, 1 sequential workflow")

agents = [
    ("🔍", "Analyst",     "Identifies at-risk students from data",         TEAL),
    ("✉️", "Communicator","Drafts personalized mentoring emails",          ACCENT),
    ("📅", "Scheduler",   "Books mentoring slots (mock calendar)",          PURPLE),
    ("📊", "Reporter",    "Writes summary report for HOD",                 NAVY),
]
box_w = Inches(2.7); gap = Inches(0.25)
total = box_w * 4 + gap * 3
sx = (SW - total) / 2
y = Inches(2.5)
for i, (ic, name, desc, col) in enumerate(agents):
    x = sx + i * (box_w + gap)
    add_round(s, x, y, box_w, Inches(3.0), fill=WHITE, line=col)
    add_rect(s, x, y, box_w, Inches(0.7), fill=col)
    add_text(s, x, y + Inches(0.13), box_w, Inches(0.5),
             name, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.9), box_w, Inches(0.8),
             ic, size=46, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(2.0), box_w - Inches(0.4), Inches(0.9),
             desc, size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    if i < 3:
        add_arrow(s, x + box_w - Inches(0.05), y + Inches(1.3),
                  gap + Inches(0.1), Inches(0.3), color=GREY_TEXT)

add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.5),
         "Same Excel as yesterday's morning demo.",
         size=16, color=GREY_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
         "But this time, the AI does the whole job.",
         size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
footer(s, 11)


# ---------- Slide 12: LIVE DEMO ----------
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_rect(s, 0, Inches(3.3), SW, Inches(0.08), fill=ACCENT)
add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
         "🎬  LIVE DEMO", size=18, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(1.4),
         "Faculty Agent Team",
         size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(3.8), Inches(12), Inches(0.6),
         "1.  Upload the same student Excel from Day 1",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.6),
         "2.  Click  ▶ Run Agent Team",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.6),
         "3.  Watch each agent's thinking, live",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.6),
         "4.  Get drafted emails, schedule, and HOD report — in 30 seconds",
         size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         "faculty_agent.py",
         size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ---------- Slide 13: Real-world examples ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Where Agentic AI is ALREADY Working",
       "Not science fiction — production today")

examples = [
    ("💻", "Coding agents",
     "Cursor, Devin, GitHub Copilot Workspace — write, run & fix code"),
    ("📞", "Customer support",
     "Klarna's AI handles 2/3 of all customer tickets autonomously"),
    ("🔬", "Research agents",
     "Deep Research (OpenAI/Google) reads 100s of sources, writes briefs"),
    ("🛒", "Shopping agents",
     "Browse + compare + checkout — Amazon's Rufus, Operator"),
    ("⚖️", "Legal agents",
     "Review contracts, flag risk clauses, draft amendments"),
    ("🏥", "Healthcare ops",
     "Schedule, triage, draft notes — Abridge, Suki"),
]
for i, (ic, t, d) in enumerate(examples):
    row = i // 3; col = i % 3
    x = Inches(0.5 + col * 4.27)
    y = Inches(1.5 + row * 2.7)
    add_round(s, x, y, Inches(4.0), Inches(2.5), fill=WHITE, line=PURPLE)
    add_text(s, x, y + Inches(0.25), Inches(4.0), Inches(0.7),
             ic, size=36, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.05), Inches(4.0), Inches(0.5),
             t, size=16, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.25), y + Inches(1.55), Inches(3.5), Inches(0.85),
             d, size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)
footer(s, 13)


# ---------- Slide 14: What this means for faculty ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "What This Means for Faculty",
       "Where to put your agent team to work")

ideas = [
    ("📥", "Admissions Triage Agent",
     "Reads applications, scores fit, drafts initial responses"),
    ("📝", "Grading Assistant Team",
     "Rubric grader + plagiarism checker + feedback writer"),
    ("📊", "Outcome Mapping Agent",
     "Maps your syllabus to NAAC / NBA criteria automatically"),
    ("🎓", "Mentor Companion",
     "Tracks each mentee, flags issues, drafts parent emails"),
    ("🔬", "Research Group Agent",
     "Reads new papers in your field, drafts a weekly digest"),
    ("📚", "Curriculum Planner",
     "Suggests course updates based on industry trends + alumni jobs"),
]
for i, (ic, t, d) in enumerate(ideas):
    row = i // 3; col = i % 3
    x = Inches(0.5 + col * 4.27)
    y = Inches(1.5 + row * 2.7)
    add_round(s, x, y, Inches(4.0), Inches(2.5), fill=WHITE, line=ACCENT)
    add_text(s, x, y + Inches(0.25), Inches(4.0), Inches(0.7),
             ic, size=36, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.05), Inches(4.0), Inches(0.5),
             t, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.25), y + Inches(1.55), Inches(3.5), Inches(0.85),
             d, size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)
footer(s, 14)


# ---------- Slide 15: Honest caveats ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Honest Caveats  —  Agents Can Fail Loudly",
       "Use them with guardrails, not blind trust")

caveats = [
    ("🔁", "They can loop forever",
     "Set max-step limits. Always have a kill switch."),
    ("💸", "They can spend money fast",
     "Every step is an LLM call. Watch tokens & set budgets."),
    ("🎯", "They can do the wrong thing",
     "Always start with HUMAN-IN-THE-LOOP — review before execution."),
    ("🛡️", "Tools = power = risk",
     "Never give an agent direct DB writes / email send without approval."),
]
for i, (ic, t, d) in enumerate(caveats):
    y = Inches(1.4 + i * 1.35)
    add_round(s, Inches(0.6), y, Inches(12.1), Inches(1.2),
              fill=WHITE, line=ACCENT)
    add_text(s, Inches(0.85), y + Inches(0.25), Inches(0.8), Inches(0.8),
             ic, size=32, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.9), y + Inches(0.15), Inches(10.5), Inches(0.5),
             t, size=18, bold=True, color=ACCENT)
    add_text(s, Inches(1.9), y + Inches(0.6), Inches(10.5), Inches(0.6),
             d, size=13, color=DARK_TEXT)
footer(s, 15)


# ---------- Slide 15.1: The Cost Elephant ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "The Elephant in the Room  —  Cost",
       "Every API call is a metered taxi ride")

add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "Quick napkin math for our Faculty Agent Team:",
         size=18, bold=True, color=NAVY)

calc = [
    ("👥",  "100 students",                "in your class"),
    ("🔄", "30 agent runs / day",         "across mentoring, grading, reports"),
    ("🧠", "~15,000 tokens / run",        "data + 4 agent prompts + responses"),
    ("📅", "200 working days / year",     ""),
    ("💸", "≈ $450 / year per faculty",   "on GPT-4o-mini alone"),
]
for i, (ic, l, r) in enumerate(calc):
    y = Inches(2.1 + i * 0.65)
    add_round(s, Inches(0.7), y, Inches(12), Inches(0.55), fill=WHITE, line=TEAL)
    add_text(s, Inches(0.9), y + Inches(0.13), Inches(0.7), Inches(0.4),
             ic, size=20, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.7), y + Inches(0.13), Inches(4.5), Inches(0.4),
             l, size=16, bold=True, color=NAVY)
    add_text(s, Inches(6.4), y + Inches(0.15), Inches(6.0), Inches(0.4),
             r, size=13, color=GREY_TEXT)

add_round(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.0),
          fill=ACCENT, line=ACCENT)
add_text(s, Inches(0.7), Inches(5.75), Inches(12), Inches(0.5),
         "Now multiply by 200 faculty in your college…",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.4),
         "$90,000+ per year. Plus all your data sitting on a US company's servers.",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
         "There IS another way →",
         size=16, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
footer(s, 16)


# ---------- Slide 15.2: Open LLMs ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Open-Source LLMs  —  Run Your Own Brain",
       "Free to download. Free to run. Free to modify. Yours.")

add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "Top open models released in 2024-2026:",
         size=18, bold=True, color=NAVY)

models = [
    ("🦙",  "Llama 3.x",       "Meta",      "70B-405B params",  "General purpose"),
    ("🌬️",  "Mistral / Mixtral","Mistral AI","7B-141B params",  "Fast, efficient"),
    ("🐉",  "Qwen 2.5",        "Alibaba",   "0.5B-72B params",  "Strong multilingual"),
    ("💎",  "Gemma 2",         "Google",    "2B-27B params",    "Lightweight, smart"),
    ("🐋",  "DeepSeek",        "DeepSeek",  "16B-671B params",  "Reasoning & code"),
    ("🔬",  "Phi-3 / Phi-4",   "Microsoft", "3B-14B params",    "Tiny but mighty"),
]
for i, (ic, name, org, size, niche) in enumerate(models):
    row = i // 3; col = i % 3
    x = Inches(0.5 + col * 4.27)
    y = Inches(2.1 + row * 2.3)
    add_round(s, x, y, Inches(4.0), Inches(2.1), fill=WHITE, line=PURPLE)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(0.8), Inches(0.6),
             ic, size=28, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.1), y + Inches(0.15), Inches(2.8), Inches(0.4),
             name, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(1.1), y + Inches(0.6), Inches(2.8), Inches(0.4),
             f"by {org}", size=11, color=GREY_TEXT)
    add_text(s, x + Inches(0.25), y + Inches(1.1), Inches(3.5), Inches(0.4),
             size, size=12, bold=True, color=PURPLE)
    add_text(s, x + Inches(0.25), y + Inches(1.5), Inches(3.5), Inches(0.5),
             niche, size=12, color=DARK_TEXT)

add_text(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
         "All of these are FREE to download from Hugging Face. Many run on a laptop.",
         size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
footer(s, 17)


# ---------- Slide 15.3: Closed vs Open tradeoffs ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Closed  vs.  Open  —  Honest Tradeoffs",
       "Neither is 'better.' They're tools for different jobs.")

# Headers
add_rect(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.7), fill=TEAL)
add_text(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.5),
         "🔒  CLOSED   (GPT, Claude, Gemini)",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.7), fill=ACCENT)
add_text(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.5),
         "🌍  OPEN   (Llama, Mistral, Qwen)",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Comparison rows
rows = [
    ("Smartness",     "Top of the leaderboard",       "Close — and catching up fast"),
    ("Cost per query","Pay per token. Adds up.",      "Free (you own the hardware)"),
    ("Data privacy",  "Leaves your campus to vendor", "Stays 100% on your servers"),
    ("Setup effort",  "Sign up. Paste key. Done.",    "Needs a GPU or decent CPU"),
    ("Customization", "Limited prompting only",       "Fine-tune on YOUR data"),
    ("Internet?",     "Always required",              "Works fully offline"),
]
for i, (label, c, o) in enumerate(rows):
    y = Inches(2.25 + i * 0.7)
    # label band
    add_rect(s, Inches(0.6), y, Inches(0.0), Inches(0.6))
    # closed side
    add_round(s, Inches(0.6), y, Inches(6.0), Inches(0.6), fill=WHITE, line=TEAL)
    add_text(s, Inches(0.8), y + Inches(0.1), Inches(2.0), Inches(0.4),
             label, size=11, bold=True, color=TEAL)
    add_text(s, Inches(2.5), y + Inches(0.13), Inches(4.0), Inches(0.4),
             c, size=12, color=DARK_TEXT)
    # open side
    add_round(s, Inches(6.8), y, Inches(6.0), Inches(0.6), fill=WHITE, line=ACCENT)
    add_text(s, Inches(7.0), y + Inches(0.13), Inches(5.6), Inches(0.4),
             o, size=12, color=DARK_TEXT)
footer(s, 18)


# ---------- Slide 15.4: Run open LLM in 30 seconds (Ollama) ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Run a Real LLM on Your Laptop  —  30 Seconds",
       "Meet Ollama. Pip-install simple, totally free.")

# 3 steps
steps = [
    ("1", "Install Ollama",        "ollama.com  →  one-click installer\n(Windows, Mac, Linux)"),
    ("2", "Pull a model",          "ollama pull llama3.2"),
    ("3", "Talk to it",            "ollama run llama3.2"),
]
for i, (n, t, d) in enumerate(steps):
    x = Inches(0.5 + i * 4.27)
    y = Inches(1.4)
    add_round(s, x, y, Inches(4.05), Inches(2.6), fill=WHITE, line=PURPLE)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.3),
                              Inches(0.8), Inches(0.8))
    circ.fill.solid(); circ.fill.fore_color.rgb = PURPLE; circ.line.fill.background()
    add_text(s, x + Inches(0.3), y + Inches(0.33), Inches(0.8), Inches(0.7),
             n, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.3), y + Inches(0.4), Inches(2.6), Inches(0.5),
             t, size=16, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), y + Inches(1.4), Inches(3.6), Inches(1.1),
             d, size=12, color=DARK_TEXT)

# Code: how to plug Ollama into our agent code
add_text(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.4),
         "And here's the punchline — switching our agents from OpenAI to Llama is a 2-line change:",
         size=13, bold=True, color=NAVY)
code = (
    "# Before — pay per token\n"
    'client = OpenAI(api_key="sk-...")\n'
    'resp = client.chat.completions.create(model="gpt-4o-mini", messages=[...])\n\n'
    "# After — runs locally, FREE\n"
    'client = OpenAI(base_url="http://localhost:11434/v1", api_key="ollama")\n'
    'resp = client.chat.completions.create(model="llama3.2", messages=[...])'
)
add_round(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.0),
          fill=NAVY, line=NAVY)
tb = add_text(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.9),
              code, size=12, color=WHITE)
for p in tb.text_frame.paragraphs:
    for r in p.runs:
        r.font.name = "Consolas"

add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
         "Same OpenAI SDK. Same prompts. Different endpoint. Zero ongoing cost.",
         size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
footer(s, 19)


# ---------- Slide 15.45a: Can a Local LLM Tell Me Today's News? ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "“Can a Local LLM Tell Me Today's News?”",
       "The #1 question faculty will ask. Here's the honest answer.")

# Big answer banner
add_round(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(1.1),
          fill=NAVY, line=NAVY)
add_text(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.5),
         "Short answer:  NO.   And that's okay.",
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(2.02), Inches(11.9), Inches(0.4),
         "A local LLM only knows what was in its training data — frozen at a cutoff date.",
         size=13, color=RGBColor(0xCF, 0xE3, 0xF5), align=PP_ALIGN.CENTER)

# Two columns: knows / doesn't know
add_round(s, Inches(0.7), Inches(2.9), Inches(5.85), Inches(3.0),
          fill=WHITE, line=TEAL)
add_rect(s, Inches(0.7), Inches(2.9), Inches(5.85), Inches(0.55), fill=TEAL)
add_text(s, Inches(0.9), Inches(2.97), Inches(5.5), Inches(0.45),
         "✅  What it CAN do", size=16, bold=True, color=WHITE)
add_bullets(s, Inches(0.9), Inches(3.55), Inches(5.5), Inches(2.3), [
    "Explain concepts, write, summarize",
    "Draft emails, lesson plans, quizzes",
    "Code & debug — fully offline",
    "Answer from its training knowledge",
], size=14)

add_round(s, Inches(6.75), Inches(2.9), Inches(5.85), Inches(3.0),
          fill=WHITE, line=ACCENT)
add_rect(s, Inches(6.75), Inches(2.9), Inches(5.85), Inches(0.55), fill=ACCENT)
add_text(s, Inches(6.95), Inches(2.97), Inches(5.5), Inches(0.45),
         "❌  What it CAN'T do", size=16, bold=True, color=WHITE)
add_bullets(s, Inches(6.95), Inches(3.55), Inches(5.5), Inches(2.3), [
    "Today's news, scores, weather",
    "Live stock / exam results / notices",
    "Read your files or the internet",
    "⚠️  If pushed — it may HALLUCINATE",
], size=14)

# Analogy
add_round(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.85),
          fill=SOFT_BLUE, line=NAVY)
add_text(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.4),
         "Analogy",
         size=12, bold=True, color=PURPLE)
add_text(s, Inches(0.9), Inches(6.42), Inches(11.5), Inches(0.45),
         "A brilliant scholar locked in a library — knows a lot, learns nothing new.",
         size=15, bold=True, color=NAVY)
footer(s, 20)


# ---------- Slide 15.45b: Brain / Memory / Hands ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Then How Do We Make It Useful?",
       "Three building blocks — and you've seen all of them in this workshop")

# Three cards
cards = [
    ("🧠", "Local LLM",
     "The BRAIN",
     "Private  •  Offline  •  Free",
     "Use for: drafting, coding, Q&A from general knowledge",
     "Demo:  local_llm_chat.py",
     TEAL),
    ("📚", "RAG",
     "+ MEMORY",
     "Grounded in YOUR documents",
     "Use for: policy docs, syllabi, lecture notes, FAQs",
     "Demo:  PolicyBot",
     PURPLE),
    ("🛠️", "Agents",
     "+ HANDS",
     "Tools — web, email, code, files",
     "Use for: live data, multi-step tasks, automation",
     "Demo:  Faculty Agent Team",
     ACCENT),
]
for i, (ic, name, role, tag, use, demo, col) in enumerate(cards):
    x = Inches(0.5 + i * 4.27)
    y = Inches(1.4)
    add_round(s, x, y, Inches(4.05), Inches(4.7), fill=WHITE, line=col)
    add_rect(s, x, y, Inches(4.05), Inches(0.95), fill=col)
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(0.8), Inches(0.7),
             ic, size=28, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.0), y + Inches(0.13), Inches(3.0), Inches(0.4),
             name, size=16, bold=True, color=WHITE)
    add_text(s, x + Inches(1.0), y + Inches(0.5), Inches(3.0), Inches(0.4),
             role, size=13, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), y + Inches(1.15), Inches(3.6), Inches(0.5),
             tag, size=13, bold=True, color=col)
    add_text(s, x + Inches(0.25), y + Inches(1.7), Inches(3.6), Inches(1.6),
             use, size=12, color=DARK_TEXT)
    add_round(s, x + Inches(0.25), y + Inches(3.7), Inches(3.55), Inches(0.85),
              fill=SOFT_BLUE, line=col, adj=0.2)
    add_text(s, x + Inches(0.35), y + Inches(3.78), Inches(3.4), Inches(0.35),
             "Workshop demo", size=10, bold=True, color=GREY_TEXT)
    add_text(s, x + Inches(0.35), y + Inches(4.05), Inches(3.4), Inches(0.45),
             demo, size=13, bold=True, color=NAVY)

# Takeaway
add_round(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.75),
          fill=NAVY, line=NAVY)
add_text(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.5),
         "Local LLM is the BRAIN.   RAG gives it MEMORY.   Agents give it HANDS.",
         size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
footer(s, 21)


# ---------- Slide 15.5: Decision matrix ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Which Should YOU Use?",
       "A practical decision guide for faculty")

choices = [
    ("🚀", "Quick prototype / demo",
     "Use CLOSED  (GPT-4o-mini, Claude Haiku)",
     "Fastest path. Pay-as-you-go. What we used in this workshop.",
     TEAL),
    ("🔒", "Student data / exam papers / IP",
     "Use OPEN  (Llama, Mistral on campus servers)",
     "Data privacy is non-negotiable. Stay local.",
     ACCENT),
    ("📈", "Production at scale (1000s/day)",
     "Use OPEN  or  HYBRID  (route easy queries to small open model, hard ones to closed)",
     "Cost flips dramatically — open wins at scale.",
     PURPLE),
    ("🎯", "Need top-tier reasoning",
     "Use CLOSED  (still ahead on the hardest tasks)",
     "GPT-4 / Claude Opus / Gemini Pro for complex reasoning.",
     TEAL),
]
for i, (ic, when, choice, why, col) in enumerate(choices):
    row = i // 2; col_i = i % 2
    x = Inches(0.5 + col_i * 6.3)
    y = Inches(1.4 + row * 2.7)
    add_round(s, x, y, Inches(6.05), Inches(2.5), fill=WHITE, line=col)
    add_rect(s, x, y, Inches(6.05), Inches(0.7), fill=col)
    add_text(s, x + Inches(0.2), y + Inches(0.13), Inches(0.7), Inches(0.5),
             ic, size=22, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(1.0), y + Inches(0.18), Inches(5.0), Inches(0.5),
             when, size=15, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), y + Inches(0.85), Inches(5.6), Inches(0.5),
             choice, size=14, bold=True, color=col)
    add_text(s, x + Inches(0.25), y + Inches(1.4), Inches(5.6), Inches(1.0),
             why, size=12, color=DARK_TEXT)
footer(s, 20)


# ---------- Slide 16: The arc of the 2 days ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "The Two-Day Arc",
       "From understanding AI  →  to commanding a team of agents")

arc = [
    ("Day 1\nSession 1",  "Intro to AI",       "Understand the brain",      TEAL),
    ("Day 1\nSession 2",  "Faculty Asst",      "First working AI app",      TEAL),
    ("Day 1\nSession 3",  "RAG / PolicyBot",   "AI grounded in your docs", TEAL),
    ("Day 2",             "Agentic AI",        "AI that ACTS",              ACCENT),
]
n = len(arc); box_w = Inches(2.8); gap = Inches(0.25)
total = box_w * n + gap * (n - 1)
sx = (SW - total) / 2; y = Inches(2.5)
for i, (tag, name, desc, col) in enumerate(arc):
    x = sx + i * (box_w + gap)
    add_round(s, x, y, box_w, Inches(2.5), fill=WHITE, line=col)
    add_rect(s, x, y, box_w, Inches(0.85), fill=col)
    add_text(s, x, y + Inches(0.12), box_w, Inches(0.7),
             tag, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.0), box_w, Inches(0.5),
             name, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(1.6), box_w - Inches(0.4), Inches(0.8),
             desc, size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    if i < n - 1:
        add_arrow(s, x + box_w - Inches(0.05), y + Inches(1.05),
                  gap + Inches(0.1), Inches(0.3), color=GREY_TEXT)

add_text(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.5),
         "Yesterday you learned how AI thinks.",
         size=18, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.5),
         "Today you saw what AI does when you give it tools.",
         size=18, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
         "Tomorrow — you build your own.",
         size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
footer(s, 16)


# ---------- Slide 17: Thank you ----------
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_rect(s, 0, Inches(3.3), SW, Inches(0.08), fill=ACCENT)
add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(1.4),
         "Thank You",
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(3.8), Inches(12), Inches(0.6),
         "You're now ahead of 99% of educators on AI.",
         size=22, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.6),
         "Go build something this week.",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
         "🎓  AI Faculty Assistant  •  📚 PolicyBot  •  🤖 Faculty Agent Team",
         size=14, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.5),
         "Questions?  Let's talk.",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)


out = "Day2_Agentic_AI.pptx"
prs.save(out)
print(f"Saved: {out}")
